"""
格式處理模組 - 統一處理文件格式轉換和保留
負責處理 Word 和 PowerPoint 之間的格式轉換邏輯
"""

from typing import Dict, Any, Optional, List
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.util import Pt
import logging


class FormatHandler:
    """格式處理器 - 負責統一處理格式轉換"""
    
    def __init__(self, logger: Optional[logging.Logger] = None):
        """
        初始化格式處理器
        
        Args:
            logger: 日誌記錄器
        """
        self.logger = logger or logging.getLogger(__name__)
    
    def extract_word_formatting(self, paragraph) -> List[Dict[str, Any]]:
        """
        從 Word 段落中提取詳細的格式信息
        
        Args:
            paragraph: Word 段落物件
            
        Returns:
            List[Dict]: 格式資訊列表
        """
        formatting_data = []
        
        try:
            for run in paragraph.runs:
                run_format = {
                    'text': run.text,
                    'font_name': None,
                    'font_size': None,
                    'font_bold': None,
                    'font_italic': None,
                    'font_underline': None,
                    'font_color': None
                }
                
                if hasattr(run, 'font'):
                    font = run.font
                    run_format.update({
                        'font_name': font.name,
                        'font_size': font.size,
                        'font_bold': font.bold,
                        'font_italic': font.italic,
                        'font_underline': font.underline
                    })
                    
                    # 提取顏色信息
                    try:
                        if hasattr(font, 'color') and font.color:
                            if hasattr(font.color, 'rgb') and font.color.rgb:
                                run_format['font_color'] = font.color.rgb
                            elif hasattr(font.color, 'theme_color') and font.color.theme_color:
                                run_format['font_color'] = 'theme_color'
                    except Exception as e:
                        self.logger.warning(f"提取字體顏色失敗: {e}")
                
                formatting_data.append(run_format)
                
        except Exception as e:
            self.logger.error(f"提取 Word 格式失敗: {e}")
            
        return formatting_data
    
    def extract_ppt_text_formatting(self, text_shape) -> Dict[str, Any]:
        """
        提取 PowerPoint 文本格式
        
        Args:
            text_shape: PowerPoint 文本形狀物件
            
        Returns:
            Dict: 格式資訊字典
        """
        formats = {'default': None, 'title': None, 'content': None}
        
        try:
            if not hasattr(text_shape, 'text_frame') or not text_shape.text_frame.paragraphs:
                return self._get_default_formats()
            
            paragraphs = text_shape.text_frame.paragraphs
            
            for i, paragraph in enumerate(paragraphs):
                if not paragraph.runs:
                    continue
                    
                first_run = paragraph.runs[0]
                format_info = self._extract_run_format(first_run, paragraph)
                
                if i == 0:
                    formats['title'] = format_info
                    if formats['default'] is None:
                        formats['default'] = format_info
                else:
                    if formats['content'] is None:
                        formats['content'] = format_info
                
                if formats['default'] is None:
                    formats['default'] = format_info
            
            # 確保有基本格式
            if formats['default'] is None:
                formats = self._get_default_formats()
            
            if formats['title'] is None:
                formats['title'] = formats['default'].copy()
                formats['title']['font_bold'] = True
            
            if formats['content'] is None:
                formats['content'] = formats['default'].copy()
                
        except Exception as e:
            self.logger.error(f"提取 PowerPoint 格式失敗: {e}")
            formats = self._get_default_formats()
        
        return formats
    
    def copy_shape_formatting(self, source_shape, target_shape) -> bool:
        """
        複製形狀格式設定
        
        Args:
            source_shape: 源形狀
            target_shape: 目標形狀
            
        Returns:
            bool: 是否成功複製
        """
        try:
            # 基本幾何屬性
            self._copy_geometry(source_shape, target_shape)
            
            # 填充格式
            self._copy_fill_format(source_shape, target_shape)
            
            # 線條格式
            self._copy_line_format(source_shape, target_shape)
            
            # 文本框架格式
            if (hasattr(source_shape, 'text_frame') and hasattr(target_shape, 'text_frame') and
                source_shape.text_frame is not None and target_shape.text_frame is not None):
                self._copy_text_frame_formatting(source_shape.text_frame, target_shape.text_frame)
            
            return True
            
        except Exception as e:
            self.logger.error(f"複製形狀格式失敗: {e}")
            return False
    
    def copy_text_frame(self, source_text_frame, target_text_frame) -> bool:
        """
        複製文本框架內容和格式
        
        Args:
            source_text_frame: 源文本框架
            target_text_frame: 目標文本框架
            
        Returns:
            bool: 是否成功複製
        """
        try:
            target_text_frame.clear()
            
            for source_paragraph in source_text_frame.paragraphs:
                # 建立目標段落
                if len(target_text_frame.paragraphs) == 1 and not target_text_frame.paragraphs[0].text:
                    target_paragraph = target_text_frame.paragraphs[0]
                else:
                    target_paragraph = target_text_frame.add_paragraph()
                
                # 複製段落格式
                self._copy_paragraph_format(source_paragraph, target_paragraph)
                
                # 複製文本運行
                for source_run in source_paragraph.runs:
                    target_run = target_paragraph.add_run()
                    target_run.text = source_run.text
                    self._copy_run_format(source_run, target_run)
            
            return True
            
        except Exception as e:
            self.logger.error(f"複製文本框架失敗: {e}")
            return False
    
    def apply_word_formatting_to_ppt(self, ppt_run, word_format: Dict[str, Any]) -> bool:
        """
        將 Word 格式應用到 PowerPoint 文本運行
        
        Args:
            ppt_run: PowerPoint 文本運行
            word_format: Word 格式資訊
            
        Returns:
            bool: 是否成功應用
        """
        try:
            if word_format.get('font_name'):
                ppt_run.font.name = word_format['font_name']
            if word_format.get('font_size'):
                ppt_run.font.size = word_format['font_size']
            if word_format.get('font_bold') is not None:
                ppt_run.font.bold = word_format['font_bold']
            if word_format.get('font_italic') is not None:
                ppt_run.font.italic = word_format['font_italic']
            if word_format.get('font_underline') is not None:
                ppt_run.font.underline = word_format['font_underline']
            
            # 應用顏色
            if word_format.get('font_color') and word_format['font_color'] != 'theme_color':
                try:
                    ppt_run.font.color.rgb = word_format['font_color']
                except Exception as e:
                    self.logger.warning(f"應用字體顏色失敗: {e}")
            
            return True
            
        except Exception as e:
            self.logger.error(f"應用 Word 格式失敗: {e}")
            return False
    
    def apply_ppt_formatting(self, run, paragraph, format_info: Dict[str, Any]) -> bool:
        """
        應用 PowerPoint 格式信息
        
        Args:
            run: 文本運行
            paragraph: 段落物件
            format_info: 格式資訊
            
        Returns:
            bool: 是否成功應用
        """
        if not format_info:
            return False
            
        try:
            # 應用字體格式
            for attr, value in format_info.items():
                if value is None:
                    continue
                    
                if attr == 'font_name':
                    run.font.name = value
                elif attr == 'font_size':
                    run.font.size = value
                elif attr == 'font_bold':
                    run.font.bold = value
                elif attr == 'font_italic':
                    run.font.italic = value
                elif attr == 'font_underline':
                    run.font.underline = value
                elif attr == 'font_color':
                    run.font.color.rgb = value
                elif attr == 'alignment':
                    paragraph.alignment = value
                elif attr == 'level':
                    paragraph.level = value
            
            return True
            
        except Exception as e:
            self.logger.error(f"應用 PowerPoint 格式失敗: {e}")
            return False
    
    # 私有方法
    def _extract_run_format(self, run, paragraph) -> Dict[str, Any]:
        """提取運行格式信息"""
        format_info = {
            'font_name': getattr(run.font, 'name', None),
            'font_size': getattr(run.font, 'size', None),
            'font_bold': getattr(run.font, 'bold', None),
            'font_italic': getattr(run.font, 'italic', None),
            'font_underline': getattr(run.font, 'underline', None),
            'font_color': None,
            'alignment': getattr(paragraph, 'alignment', None),
            'level': getattr(paragraph, 'level', None)
        }
        
        try:
            if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'):
                format_info['font_color'] = run.font.color.rgb
        except:
            pass
        
        return format_info
    
    def _get_default_formats(self) -> Dict[str, Any]:
        """獲取預設格式"""
        basic_format = {
            'font_name': 'Arial',
            'font_size': Pt(24),
            'font_bold': False,
            'font_italic': False,
            'font_underline': None,
            'font_color': None,
            'alignment': None,
            'level': 0
        }
        
        return {
            'default': basic_format,
            'title': {**basic_format, 'font_bold': True},
            'content': basic_format.copy()
        }
    
    def _copy_geometry(self, source_shape, target_shape):
        """複製幾何屬性"""
        target_shape.left = source_shape.left
        target_shape.top = source_shape.top
        target_shape.width = source_shape.width
        target_shape.height = source_shape.height
    
    def _copy_fill_format(self, source_shape, target_shape):
        """複製填充格式"""
        if not (hasattr(source_shape, 'fill') and hasattr(target_shape, 'fill')):
            return
            
        try:
            if hasattr(source_shape.fill, 'type'):
                fill_type = source_shape.fill.type
                if fill_type == MSO_FILL_TYPE.SOLID:
                    target_shape.fill.solid()
                    if (hasattr(source_shape.fill, 'fore_color') and 
                        hasattr(source_shape.fill.fore_color, 'rgb')):
                        target_shape.fill.fore_color.rgb = source_shape.fill.fore_color.rgb
                else:
                    target_shape.fill.background()
            else:
                target_shape.fill.background()
        except Exception as e:
            self.logger.warning(f"複製填充格式失敗: {e}")
            try:
                target_shape.fill.background()
            except:
                pass
    
    def _copy_line_format(self, source_shape, target_shape):
        """複製線條格式"""
        if not (hasattr(source_shape, 'line') and hasattr(target_shape, 'line')):
            return
            
        try:
            line_width_pt = 0
            if (hasattr(source_shape.line, 'width') and 
                source_shape.line.width is not None and 
                hasattr(source_shape.line.width, 'pt')):
                line_width_pt = source_shape.line.width.pt
            
            if line_width_pt > 0:
                target_shape.line.width = source_shape.line.width
                if (hasattr(source_shape.line, 'color') and hasattr(target_shape.line, 'color') and
                    hasattr(source_shape.line.color, 'rgb')):
                    target_shape.line.color.rgb = source_shape.line.color.rgb
            else:
                target_shape.line.width = Pt(0)
        except Exception as e:
            self.logger.warning(f"複製線條格式失敗: {e}")
    
    def _copy_paragraph_format(self, source_paragraph, target_paragraph):
        """複製段落格式"""
        try:
            if hasattr(source_paragraph, 'alignment'):
                target_paragraph.alignment = source_paragraph.alignment
            if hasattr(source_paragraph, 'level'):
                target_paragraph.level = source_paragraph.level
        except Exception as e:
            self.logger.warning(f"複製段落格式失敗: {e}")
    
    def _copy_run_format(self, source_run, target_run):
        """複製運行格式"""
        try:
            if hasattr(source_run, 'font') and hasattr(target_run, 'font'):
                source_font = source_run.font
                target_font = target_run.font
                
                if source_font.name:
                    target_font.name = source_font.name
                if source_font.size:
                    target_font.size = source_font.size
                if source_font.bold is not None:
                    target_font.bold = source_font.bold
                if source_font.italic is not None:
                    target_font.italic = source_font.italic
                if source_font.underline is not None:
                    target_font.underline = source_font.underline
                if hasattr(source_font, 'color') and source_font.color:
                    target_font.color.rgb = source_font.color.rgb
        except Exception as e:
            self.logger.warning(f"複製運行格式失敗: {e}")
    
    def _copy_text_frame_formatting(self, source_text_frame, target_text_frame):
        """複製文本框架格式"""
        try:
            # 複製邊距設置
            for margin in ['margin_left', 'margin_right', 'margin_top', 'margin_bottom']:
                if hasattr(source_text_frame, margin) and hasattr(target_text_frame, margin):
                    setattr(target_text_frame, margin, getattr(source_text_frame, margin))
            
            # 複製其他格式設置
            for attr in ['auto_size', 'vertical_anchor', 'word_wrap']:
                if hasattr(source_text_frame, attr) and hasattr(target_text_frame, attr):
                    setattr(target_text_frame, attr, getattr(source_text_frame, attr))
                    
        except Exception as e:
            self.logger.warning(f"複製文本框架格式失敗: {e}")