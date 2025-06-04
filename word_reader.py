"""
Word 文档读取工具
使用 python-docx 库来读取 Word 文档内容
"""

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE, MSO_LINE_DASH_STYLE
from typing import List, Dict, Optional
import os
import re
import io


def read_word_document(file_path: str) -> Dict[str, any]:
    """
    读取 Word 文档内容
    
    Args:
        file_path (str): Word 文档的文件路径
        
    Returns:
        Dict[str, any]: 包含文档内容的字典，包括:
            - text: 纯文本内容
            - paragraphs: 段落列表
            - tables: 表格内容（如果有）
            - metadata: 文档元数据
            - success: 操作是否成功
            - error: 错误信息（如果有）
    """
    result = {
        'text': '',
        'paragraphs': [],
        'tables': [],
        'metadata': {},
        'success': False,
        'error': None
    }
    
    try:
        # 检查文件是否存在
        if not os.path.exists(file_path):
            result['error'] = f"文件不存在: {file_path}"
            return result
        
        # 检查文件扩展名
        if not file_path.lower().endswith('.docx'):
            result['error'] = "文件必须是 .docx 格式"
            return result
        
        # 读取文档
        doc = Document(file_path)
        
        # 提取段落内容
        paragraphs = []
        full_text = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():  # 忽略空段落
                paragraphs.append({
                    'text': paragraph.text,
                    'style': paragraph.style.name if paragraph.style else None
                })
                full_text.append(paragraph.text)
        
        result['paragraphs'] = paragraphs
        result['text'] = '\n'.join(full_text)
        
        # 提取表格内容
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            tables.append(table_data)
        
        result['tables'] = tables
        
        # 提取文档元数据
        core_props = doc.core_properties
        metadata = {
            'title': core_props.title,
            'author': core_props.author,
            'subject': core_props.subject,
            'created': core_props.created,
            'modified': core_props.modified,
            'category': core_props.category,
            'comments': core_props.comments
        }
        result['metadata'] = metadata
        
        result['success'] = True
        
    except Exception as e:
        result['error'] = f"读取文档时发生错误: {str(e)}"
    
    return result


def read_powerpoint_document(file_path: str) -> Dict[str, any]:
    """
    读取 PowerPoint 文档内容
    
    Args:
        file_path (str): PowerPoint 文档的文件路径
        
    Returns:
        Dict[str, any]: 包含文档内容的字典，包括:
            - slides: 幻灯片列表
            - total_slides: 总幻灯片数
            - text: 所有文本内容
            - success: 操作是否成功
            - error: 错误信息（如果有）
    """
    result = {
        'slides': [],
        'total_slides': 0,
        'text': '',
        'success': False,
        'error': None
    }
    
    try:
        # 检查文件是否存在
        if not os.path.exists(file_path):
            result['error'] = f"文件不存在: {file_path}"
            return result
        
        # 检查文件扩展名
        if not file_path.lower().endswith('.pptx'):
            result['error'] = "文件必须是 .pptx 格式"
            return result
        
        # 读取演示文稿
        prs = Presentation(file_path)
        
        slides = []
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                'slide_number': slide_num,
                'title': '',
                'content': [],
                'text_runs': [],
                'layout_name': slide.slide_layout.name if slide.slide_layout else 'Unknown'
            }
            
            # 提取幻灯片中的所有文本
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content = shape.text.strip()
                    slide_text.append(text_content)
                    slide_data['text_runs'].append({
                        'text': text_content,
                        'shape_type': str(shape.shape_type) if hasattr(shape, 'shape_type') else 'Unknown'
                    })
                    
                    # 尝试识别标题
                    if not slide_data['title'] and (
                        len(text_content) < 100 or
                        'title' in text_content.lower()[:20]
                    ):
                        slide_data['title'] = text_content
                
                # 处理表格内容
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_data.append(cell_text)
                                slide_text.append(cell_text)
                        if row_data:
                            table_data.append(row_data)
                    
                    if table_data:
                        slide_data['text_runs'].append({
                            'text': f"[表格: {len(table_data)}行]",
                            'shape_type': 'Table',
                            'table_data': table_data
                        })
            
            slide_data['content'] = slide_text
            slide_data['full_text'] = '\n'.join(slide_text)
            
            # 如果没有识别到标题，使用第一行文本作为标题
            if not slide_data['title'] and slide_text:
                slide_data['title'] = slide_text[0][:50] + ('...' if len(slide_text[0]) > 50 else '')
            
            slides.append(slide_data)
            all_text.extend(slide_text)
        
        result['slides'] = slides
        result['total_slides'] = len(slides)
        result['text'] = '\n'.join(all_text)
        result['success'] = True
        
    except Exception as e:
        result['error'] = f"读取PowerPoint文档时发生错误: {str(e)}"
    
    return result


def get_word_text_only(file_path: str) -> str:
    """
    仅获取 Word 文档的纯文本内容
    
    Args:
        file_path (str): Word 文档的文件路径
        
    Returns:
        str: 文档的纯文本内容，如果出错则返回空字符串
    """
    result = read_word_document(file_path)
    if result['success']:
        return result['text']
    else:
        print(f"读取文档失败: {result['error']}")
        return ""


def get_word_paragraphs(file_path: str) -> List[str]:
    """
    获取 Word 文档的段落列表
    
    Args:
        file_path (str): Word 文档的文件路径
        
    Returns:
        List[str]: 段落文本列表
    """
    result = read_word_document(file_path)
    if result['success']:
        return [p['text'] for p in result['paragraphs']]
    else:
        print(f"读取文档失败: {result['error']}")
        return []


def print_word_summary(file_path: str) -> None:
    """
    打印 Word 文档的摘要信息
    
    Args:
        file_path (str): Word 文档的文件路径
    """
    result = read_word_document(file_path)
    
    if not result['success']:
        print(f"❌ 读取失败: {result['error']}")
        return
    
    print("📄 Word 文档摘要")
    print("=" * 50)
    print(f"文件路径: {file_path}")
    print(f"段落数量: {len(result['paragraphs'])}")
    print(f"表格数量: {len(result['tables'])}")
    print(f"字符数量: {len(result['text'])}")
    
    # 显示元数据
    metadata = result['metadata']
    if any(metadata.values()):
        print("\n📊 文档元数据:")
        for key, value in metadata.items():
            if value:
                print(f"  {key}: {value}")
    
    # 显示所有段落
    if result['paragraphs']:
        print(f"\n📝 所有段落内容:")
        for i, paragraph in enumerate(result['paragraphs']):
            print(f"  {i+1}. {paragraph['text']}")
    
    # 显示表格信息
    if result['tables']:
        print(f"\n📋 表格信息:")
        for i, table in enumerate(result['tables']):
            print(f"  表格 {i+1}: {len(table)} 行 x {len(table[0]) if table else 0} 列")


def parse_numbered_sections(file_path: str) -> Dict[str, any]:
    """
    将 Word 文档内容按照 "{数字}." 格式进行分段
    
    Args:
        file_path (str): Word 文档的文件路径
        
    Returns:
        Dict[str, any]: 包含分段内容的字典，包括:
            - sections: 分段列表，每个段包含编号和内容
            - total_sections: 总段数
            - success: 操作是否成功
            - error: 错误信息（如果有）
    """
    
    result = {
        'sections': [],
        'total_sections': 0,
        'success': False,
        'error': None
    }
    
    try:
        # 先读取文档内容
        doc_result = read_word_document(file_path)
        if not doc_result['success']:
            result['error'] = doc_result['error']
            return result
        
        paragraphs = doc_result['paragraphs']
        sections = []
        current_section = None
        
        # 正则表达式匹配 "数字." 格式
        number_pattern = re.compile(r'^(\d+)\.\s*(.*)')
        
        for paragraph in paragraphs:
            text = paragraph['text'].strip()
            if not text:
                continue
            
            # 检查是否是新段落开始（以数字.开头）
            match = number_pattern.match(text)
            if match:
                # 如果有当前段落，先保存它
                if current_section is not None:
                    sections.append(current_section)
                
                # 开始新段落
                section_number = int(match.group(1))
                section_title = match.group(2) if match.group(2) else ""
                
                current_section = {
                    'number': section_number,
                    'title': section_title,
                    'content': [text],  # 包含标题行
                    'text_only': section_title  # 仅文本内容（不含编号）
                }
            else:
                # 如果不是新段落开始，添加到当前段落
                if current_section is not None:
                    current_section['content'].append(text)
                    if current_section['text_only']:
                        current_section['text_only'] += '\n' + text
                    else:
                        current_section['text_only'] = text
                else:
                    # 如果还没有开始任何段落，创建一个"前言"段落
                    current_section = {
                        'number': 0,
                        'title': '前言',
                        'content': [text],
                        'text_only': text
                    }
        
        # 保存最后一个段落
        if current_section is not None:
            sections.append(current_section)
        
        result['sections'] = sections
        result['total_sections'] = len(sections)
        result['success'] = True
        
    except Exception as e:
        result['error'] = f"分段处理时发生错误: {str(e)}"
    
    return result


def get_section_by_number(file_path: str, section_number: int) -> Dict[str, any]:
    """
    根据编号获取特定段落
    
    Args:
        file_path (str): Word 文档的文件路径
        section_number (int): 段落编号
        
    Returns:
        Dict[str, any]: 包含段落信息的字典
    """
    sections_result = parse_numbered_sections(file_path)
    
    if not sections_result['success']:
        return {
            'section': None,
            'success': False,
            'error': sections_result['error']
        }
    
    # 查找指定编号的段落
    target_section = None
    for section in sections_result['sections']:
        if section['number'] == section_number:
            target_section = section
            break
    
    return {
        'section': target_section,
        'success': target_section is not None,
        'error': f"未找到编号为 {section_number} 的段落" if target_section is None else None
    }


def print_sections_summary(file_path: str) -> None:
    """
    打印文档分段摘要
    
    Args:
        file_path (str): Word 文档的文件路径
    """
    result = parse_numbered_sections(file_path)
    
    if not result['success']:
        print(f"❌ 分段失败: {result['error']}")
        return
    
    print("📚 文档分段摘要")
    print("=" * 60)
    print(f"总段数: {result['total_sections']}")
    print("\n📋 段落列表:")
    
    for section in result['sections']:
        title_preview = section['title'][:50] + "..." if len(section['title']) > 50 else section['title']
        content_lines = len(section['content'])
        print(f"  {section['number']:3d}. {title_preview} ({content_lines} 行)")
    
    print(f"\n💡 使用 get_section_by_number(file_path, 编号) 可以获取特定段落")


def print_section_detail(file_path: str, section_number: int) -> None:
    """
    打印特定段落的详细内容
    
    Args:
        file_path (str): Word 文档的文件路径
        section_number (int): 段落编号
    """
    result = get_section_by_number(file_path, section_number)
    
    if not result['success']:
        print(f"❌ {result['error']}")
        return
    
    section = result['section']
    print(f"📖 段落 {section['number']} 详细内容")
    print("=" * 60)
    print(f"标题: {section['title']}")
    print(f"行数: {len(section['content'])}")
    print("\n📝 完整内容:")
    print("-" * 40)
    for line in section['content']:
        print(line)
    print("-" * 40)


def print_powerpoint_summary(file_path: str) -> None:
    """
    打印 PowerPoint 文档的摘要信息
    
    Args:
        file_path (str): PowerPoint 文档的文件路径
    """
    result = read_powerpoint_document(file_path)
    
    if not result['success']:
        print(f"❌ 读取失败: {result['error']}")
        return
    
    print("🎯 PowerPoint 文档摘要")
    print("=" * 60)
    print(f"文件路径: {file_path}")
    print(f"幻灯片数量: {result['total_slides']}")
    print(f"总字符数: {len(result['text'])}")
    
    print(f"\n📋 幻灯片列表:")
    for slide in result['slides']:
        title_preview = slide['title'][:60] + "..." if len(slide['title']) > 60 else slide['title']
        content_count = len(slide['content'])
        layout = slide['layout_name']
        print(f"  第{slide['slide_number']:2d}页: {title_preview}")
        print(f"      └─ 布局: {layout}, 内容块: {content_count}个")


def print_slide_detail(file_path: str, slide_number: int) -> None:
    """
    打印特定幻灯片的详细内容
    
    Args:
        file_path (str): PowerPoint 文档的文件路径
        slide_number (int): 幻灯片编号（从1开始）
    """
    result = read_powerpoint_document(file_path)
    
    if not result['success']:
        print(f"❌ 读取失败: {result['error']}")
        return
    
    if slide_number < 1 or slide_number > result['total_slides']:
        print(f"❌ 幻灯片编号无效，请输入1到{result['total_slides']}之间的数字")
        return
    
    slide = result['slides'][slide_number - 1]
    print(f"🎯 第{slide_number}页幻灯片详细内容")
    print("=" * 60)
    print(f"标题: {slide['title']}")
    print(f"布局: {slide['layout_name']}")
    print(f"内容块数: {len(slide['text_runs'])}")
    
    print("\n📝 详细内容:")
    print("-" * 40)
    for i, text_run in enumerate(slide['text_runs'], 1):
        print(f"[{i}] {text_run['shape_type']}: {text_run['text']}")
        if 'table_data' in text_run:
            print("    表格内容:")
            for row_idx, row in enumerate(text_run['table_data']):
                print(f"      第{row_idx+1}行: {' | '.join(row)}")
    print("-" * 40)


def duplicate_slide(file_path: str, source_slide_number: int, copy_count: int = 1, output_file: str = None) -> Dict[str, any]:
    """
    复制指定的幻灯片，支持保存到新文件
    
    Args:
        file_path (str): PowerPoint 文档的文件路径
        source_slide_number (int): 源幻灯片编号（从1开始）
        copy_count (int): 复制的数量，默认为1
        output_file (str): 输出文件路径，如果为None则覆盖原文件
        
    Returns:
        Dict[str, any]: 包含操作结果的字典
    """
    result = {
        'success': False,
        'error': None,
        'copied_slides': 0,
        'total_slides_after': 0,
        'output_file': output_file or file_path
    }
    
    try:
        # 检查文件是否存在
        if not os.path.exists(file_path):
            result['error'] = f"文件不存在: {file_path}"
            return result
        
        # 读取演示文稿
        prs = Presentation(file_path)
        
        # 验证源幻灯片编号
        if source_slide_number < 1 or source_slide_number > len(prs.slides):
            result['error'] = f"源幻灯片编号无效，请输入1到{len(prs.slides)}之间的数字"
            return result
        
        # 获取源幻灯片（注意：索引从0开始）
        source_slide = prs.slides[source_slide_number - 1]
        source_layout = source_slide.slide_layout
        
        # 复制指定次数
        for i in range(copy_count):
            # 创建新幻灯片，使用相同的布局
            new_slide = prs.slides.add_slide(source_layout)
            
            # 完整复制幻灯片内容，包括背景
            copy_slide_completely(source_slide, new_slide)
        
        # 保存文件（保存到指定的输出文件或原文件）
        save_path = output_file if output_file else file_path
        prs.save(save_path)
        
        result['success'] = True
        result['copied_slides'] = copy_count
        result['total_slides_after'] = len(prs.slides)
        result['output_file'] = save_path
        
    except Exception as e:
        result['error'] = f"复制幻灯片时发生错误: {str(e)}"
    
    return result


def copy_slide_completely(source_slide, target_slide):
    """
    完整复制幻灯片内容，包括背景、形状和格式
    
    Args:
        source_slide: 源幻灯片
        target_slide: 目标幻灯片
    """
    try:
        # 1. 复制幻灯片背景
        copy_slide_background(source_slide, target_slide)
        
        # 2. 复制所有形状
        for shape in source_slide.shapes:
            try:
                if shape.is_placeholder:
                    # 处理占位符
                    copy_placeholder_content(shape, target_slide)
                else:
                    # 处理非占位符形状
                    copy_non_placeholder_shape(shape, target_slide)
            except Exception as e:
                print(f"复制形状时出错: {str(e)}")
                continue
        
    except Exception as e:
        print(f"完整复制幻灯片时出错: {str(e)}")


def copy_slide_background(source_slide, target_slide):
    """
    複製幻燈片背景，包括背景圖片
    
    Args:
        source_slide: 源幻灯片
        target_slide: 目标幻灯片
    """
    try:
        # 方法1: 嘗試複製背景填充
        if hasattr(source_slide, 'background'):
            try:
                source_bg = source_slide.background
                target_bg = target_slide.background
                
                # 檢查背景填充類型
                if hasattr(source_bg, 'fill') and hasattr(target_bg, 'fill'):
                    source_fill = source_bg.fill
                    target_fill = target_bg.fill
                    
                    # 複製填充類型和屬性
                    if hasattr(source_fill, 'type'):
                        fill_type = source_fill.type
                        
                        if fill_type == MSO_FILL_TYPE.SOLID:
                            # 純色背景
                            target_fill.solid()
                            if hasattr(source_fill, 'fore_color'):
                                target_fill.fore_color.rgb = source_fill.fore_color.rgb
                        elif fill_type == MSO_FILL_TYPE.PICTURE:
                            # 圖片背景 - 這是我們最關心的
                            try:
                                # 嘗試複製圖片背景
                                if hasattr(source_fill, 'fore_color') and hasattr(source_fill.fore_color, 'rgb'):
                                    target_fill.solid()
                                    target_fill.fore_color.rgb = source_fill.fore_color.rgb
                                print("   📸 嘗試複製背景圖片...")
                            except Exception as pic_error:
                                print(f"   ⚠️  複製背景圖片時出錯: {str(pic_error)}")
                        elif fill_type == MSO_FILL_TYPE.GRADIENT:
                            # 漸變背景
                            try:
                                if hasattr(source_fill, 'gradient') and hasattr(target_fill, 'gradient'):
                                    # 複製漸變設置（簡化版）
                                    target_fill.gradient()
                                    print("   🎨 嘗試複製漸變背景...")
                            except:
                                pass
                        else:
                            # 其他類型背景，設為透明
                            target_fill.background()
                            
            except Exception as bg_error:
                print(f"   ⚠️  複製背景填充時出錯: {str(bg_error)}")
        
        # 方法2: 檢查幻燈片母版和佈局
        try:
            # 確保使用相同的佈局（這有助於保持背景一致性）
            if hasattr(source_slide, 'slide_layout') and hasattr(target_slide, 'slide_layout'):
                source_layout = source_slide.slide_layout
                target_layout = target_slide.slide_layout
                
                # 如果佈局不同，記錄差異
                if source_layout.name != target_layout.name:
                    print(f"   ⚠️  佈局差異: 源={source_layout.name}, 目標={target_layout.name}")
                    
        except Exception as layout_error:
            print(f"   ⚠️  檢查佈局時出錯: {str(layout_error)}")
        
        # 方法3: 檢查是否有背景形狀（有時背景是作為形狀存在的）
        try:
            # 查找可能的背景形狀（通常在最底層）
            background_shapes = []
            for shape in source_slide.shapes:
                # 檢查是否是可能的背景形狀（大尺寸、在底層的圖片或矩形）
                if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE or 
                    shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE):
                    # 檢查形狀大小（如果接近幻燈片大小，可能是背景）
                    slide_width = source_slide.slide_layout.width if hasattr(source_slide.slide_layout, 'width') else 0
                    slide_height = source_slide.slide_layout.height if hasattr(source_slide.slide_layout, 'height') else 0
                    
                    if (slide_width > 0 and slide_height > 0 and
                        shape.width >= slide_width * 0.8 and  # 寬度達到幻燈片的80%以上
                        shape.height >= slide_height * 0.8):  # 高度達到幻燈片的80%以上
                        background_shapes.append(shape)
            
            if background_shapes:
                print(f"   🖼️  發現 {len(background_shapes)} 個可能的背景形狀")
                
        except Exception as shape_error:
            print(f"   ⚠️  檢查背景形狀時出錯: {str(shape_error)}")
            
    except Exception as e:
        print(f"   ❌ 複製幻燈片背景時出錯: {str(e)}")


def copy_placeholder_content(source_placeholder, target_slide):
    """
    复制占位符内容
    
    Args:
        source_placeholder: 源占位符
        target_slide: 目标幻灯片
    """
    try:
        # 找到目标幻灯片中对应的占位符
        for target_shape in target_slide.shapes:
            if (target_shape.is_placeholder and 
                hasattr(source_placeholder, 'placeholder_format') and
                hasattr(target_shape, 'placeholder_format') and
                target_shape.placeholder_format.idx == source_placeholder.placeholder_format.idx):
                
                # 复制文本内容
                if hasattr(source_placeholder, 'text') and hasattr(target_shape, 'text'):
                    target_shape.text = source_placeholder.text
                
                # 复制文本框架内容（包括格式）
                if hasattr(source_placeholder, 'text_frame') and hasattr(target_shape, 'text_frame'):
                    copy_text_frame(source_placeholder.text_frame, target_shape.text_frame)
                
                # 复制表格内容
                if hasattr(source_placeholder, 'table') and source_placeholder.has_table:
                    copy_table_content(source_placeholder.table, target_shape)
                
                break
    except Exception as e:
        print(f"复制占位符内容时出错: {str(e)}")


def copy_non_placeholder_shape(source_shape, target_slide):
    """
    复制非占位符形状
    
    Args:
        source_shape: 源形状
        target_slide: 目标幻灯片
    """
    try:
        from pptx.util import Inches
        
        # 获取形状的基本属性
        left = source_shape.left
        top = source_shape.top
        width = source_shape.width
        height = source_shape.height
        
        if source_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # 复制文本框
            new_textbox = target_slide.shapes.add_textbox(left, top, width, height)
            if hasattr(source_shape, 'text_frame'):
                copy_text_frame(source_shape.text_frame, new_textbox.text_frame)
            copy_shape_formatting(source_shape, new_textbox)
            
        elif source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # 复制图片
            try:
                # 获取图片数据
                image_blob = source_shape.image.blob
                image_stream = io.BytesIO(image_blob)
                new_picture = target_slide.shapes.add_picture(image_stream, left, top, width, height)
                copy_shape_formatting(source_shape, new_picture)
            except Exception as e:
                print(f"复制图片时出错: {str(e)}")
                
        elif source_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # 复制自动形状
            try:
                if hasattr(source_shape, 'auto_shape_type'):
                    new_shape = target_slide.shapes.add_shape(
                        source_shape.auto_shape_type, left, top, width, height
                    )
                    if hasattr(source_shape, 'text_frame'):
                        copy_text_frame(source_shape.text_frame, new_shape.text_frame)
                    copy_shape_formatting(source_shape, new_shape)
            except Exception as e:
                print(f"复制自动形状时出错: {str(e)}")
                
        elif source_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # 复制组合形状（简化处理）
            print(f"检测到组合形状，跳过复制")
            
        elif source_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # 复制表格
            try:
                if hasattr(source_shape, 'table'):
                    table = source_shape.table
                    rows = len(table.rows)
                    cols = len(table.columns)
                    new_table = target_slide.shapes.add_table(rows, cols, left, top, width, height)
                    copy_table_content(table, new_table)
            except Exception as e:
                print(f"复制表格时出错: {str(e)}")
        
        else:
            # 其他类型的形状，尝试基本复制
            print(f"检测到未处理的形状类型: {source_shape.shape_type}")
            
    except Exception as e:
        print(f"复制非占位符形状时出错: {str(e)}")


def copy_shape_formatting(source_shape, target_shape):
    """
    复制形状的格式设置，精确保留原始格式
    
    Args:
        source_shape: 源形状
        target_shape: 目标形状
    """
    try:
        # 1. 复制基本几何属性（位置和大小）
        try:
            target_shape.left = source_shape.left
            target_shape.top = source_shape.top  
            target_shape.width = source_shape.width
            target_shape.height = source_shape.height
        except Exception as geom_error:
            print(f"     ⚠️  複製幾何屬性時出錯: {str(geom_error)}")
        
        # 2. 复制填充格式
        if hasattr(source_shape, 'fill') and hasattr(target_shape, 'fill'):
            try:
                # 检查源形状的填充类型
                if hasattr(source_shape.fill, 'type'):
                    fill_type = source_shape.fill.type
                    
                    if fill_type == MSO_FILL_TYPE.SOLID:
                        # 纯色填充
                        target_shape.fill.solid()
                        if hasattr(source_shape.fill, 'fore_color') and hasattr(source_shape.fill.fore_color, 'rgb'):
                            target_shape.fill.fore_color.rgb = source_shape.fill.fore_color.rgb
                            print(f"     ✅ 複製純色填充: {source_shape.fill.fore_color.rgb}")
                    elif fill_type == MSO_FILL_TYPE.BACKGROUND:
                        # 背景填充（通常用于透明）
                        target_shape.fill.background()
                        print(f"     ✅ 設置背景填充（透明）")
                    elif fill_type == MSO_FILL_TYPE.GRADIENT:
                        # 渐变填充 - 嘗試複製基本漸變
                        try:
                            target_shape.fill.gradient()
                            print(f"     ✅ 設置漸變填充")
                        except:
                            target_shape.fill.background()
                            print(f"     ⚠️  漸變複製失敗，設為透明")
                    elif fill_type == MSO_FILL_TYPE.PATTERN:
                        # 图案填充 - 暂时设为背景色
                        target_shape.fill.background()
                        print(f"     ⚠️  圖案填充暫不支持，設為透明")
                    elif fill_type == MSO_FILL_TYPE.PICTURE:
                        # 图片填充 - 暂时设为背景色
                        target_shape.fill.background()
                        print(f"     ⚠️  圖片填充暫不支持，設為透明")
                    else:
                        # 其他类型或未知类型，设为背景色（透明）
                        target_shape.fill.background()
                        print(f"     ⚠️  未知填充類型: {fill_type}，設為透明")
                else:
                    # 无法确定填充类型，设为背景色（透明）
                    target_shape.fill.background()
                    print(f"     ⚠️  無法確定填充類型，設為透明")
            except Exception as fill_error:
                print(f"     ❌ 复制填充格式时出错: {str(fill_error)}")
                # 出错时设为透明背景
                try:
                    target_shape.fill.background()
                except:
                    pass
        
        # 3. 复制线条格式 - 重要：確保正確處理無邊框情況
        if hasattr(source_shape, 'line') and hasattr(target_shape, 'line'):
            try:
                # 首先檢查源形狀是否有可見的線條
                source_has_visible_line = False
                line_width_pt = 0
                
                try:
                    # 檢查線條寬度是否大於0
                    if (hasattr(source_shape.line, 'width') and 
                        source_shape.line.width is not None and 
                        hasattr(source_shape.line.width, 'pt')):
                        line_width_pt = source_shape.line.width.pt
                        if line_width_pt > 0:
                            source_has_visible_line = True
                except:
                    pass
                
                if source_has_visible_line:
                    # 源形狀有可見線條，複製線條屬性
                    print(f"     📐 複製線條格式，寬度: {line_width_pt}pt")
                    
                    # 複製線條寬度
                    if hasattr(source_shape.line, 'width') and hasattr(target_shape.line, 'width'):
                        try:
                            target_shape.line.width = source_shape.line.width
                        except Exception as width_error:
                            print(f"     ⚠️  複製線條寬度失敗: {str(width_error)}")
                    
                    # 複製線條顏色
                    if hasattr(source_shape.line, 'color') and hasattr(target_shape.line, 'color'):
                        try:
                            if hasattr(source_shape.line.color, 'rgb'):
                                target_shape.line.color.rgb = source_shape.line.color.rgb
                                print(f"     ✅ 複製線條顏色: {source_shape.line.color.rgb}")
                        except Exception as color_error:
                            print(f"     ⚠️  複製線條顏色失敗: {str(color_error)}")
                    
                    # 複製線條樣式
                    if hasattr(source_shape.line, 'dash_style') and hasattr(target_shape.line, 'dash_style'):
                        try:
                            target_shape.line.dash_style = source_shape.line.dash_style
                        except Exception as dash_error:
                            print(f"     ⚠️  複製線條樣式失敗: {str(dash_error)}")
                else:
                    # 源形狀沒有可見線條，移除目標形狀的線條
                    try:
                        from pptx.util import Pt
                        target_shape.line.width = Pt(0)  # 設置線條寬度為0
                        print(f"     ✅ 移除目標形狀線條（設為無邊框）")
                    except Exception as remove_error:
                        print(f"     ⚠️  移除線條失敗: {str(remove_error)}")
                    
            except Exception as line_error:
                print(f"     ❌ 复制线条格式时出错: {str(line_error)}")
                # 出錯時也嘗試移除線條
                try:
                    from pptx.util import Pt
                    target_shape.line.width = Pt(0)
                except:
                    pass
        
        # 4. 复制文本框架的特殊格式（如果是文本框）
        if (hasattr(source_shape, 'text_frame') and hasattr(target_shape, 'text_frame') and
            source_shape.text_frame is not None and target_shape.text_frame is not None):
            try:
                copy_text_frame_formatting(source_shape.text_frame, target_shape.text_frame)
                print(f"     ✅ 複製文本框架格式")
            except Exception as text_frame_error:
                print(f"     ⚠️  複製文本框架格式失敗: {str(text_frame_error)}")
        
        # 5. 复制阴影效果
        if hasattr(source_shape, 'shadow') and hasattr(target_shape, 'shadow'):
            try:
                # 檢查是否有陰影
                if hasattr(source_shape.shadow, 'inherit'):
                    target_shape.shadow.inherit = source_shape.shadow.inherit
                    print(f"     ✅ 複製陰影設置")
            except Exception as shadow_error:
                print(f"     ⚠️  複製陰影失敗: {str(shadow_error)}")
        
        # 6. 复制旋转角度
        try:
            if hasattr(source_shape, 'rotation') and hasattr(target_shape, 'rotation'):
                target_shape.rotation = source_shape.rotation
                if source_shape.rotation != 0:
                    print(f"     ✅ 複製旋轉角度: {source_shape.rotation}°")
        except Exception as rotation_error:
            print(f"     ⚠️  複製旋轉角度失敗: {str(rotation_error)}")
                
    except Exception as e:
        print(f"     ❌ 复制形状格式时出错: {str(e)}")


def copy_table_content(source_table, target_table_shape):
    """
    复制表格内容
    
    Args:
        source_table: 源表格
        target_table_shape: 目标表格形状或表格对象
    """
    try:
        # 确保目标是表格对象
        if hasattr(target_table_shape, 'table'):
            target_table = target_table_shape.table
        else:
            target_table = target_table_shape
        
        # 复制表格内容
        for row_idx, source_row in enumerate(source_table.rows):
            if row_idx < len(target_table.rows):
                target_row = target_table.rows[row_idx]
                for col_idx, source_cell in enumerate(source_row.cells):
                    if col_idx < len(target_row.cells):
                        target_cell = target_row.cells[col_idx]
                        target_cell.text = source_cell.text
                        
                        # 复制单元格文本格式
                        if hasattr(source_cell, 'text_frame') and hasattr(target_cell, 'text_frame'):
                            copy_text_frame(source_cell.text_frame, target_cell.text_frame)
    except Exception as e:
        print(f"复制表格内容时出错: {str(e)}")


def copy_text_frame(source_text_frame, target_text_frame):
    """
    复制文本框架的内容和格式
    
    Args:
        source_text_frame: 源文本框架
        target_text_frame: 目标文本框架
    """
    try:
        # 清空目标文本框架
        target_text_frame.clear()
        
        # 复制每个段落
        for source_paragraph in source_text_frame.paragraphs:
            # 在目标中添加新段落（除了第一个，它已经存在）
            if len(target_text_frame.paragraphs) == 1 and not target_text_frame.paragraphs[0].text:
                target_paragraph = target_text_frame.paragraphs[0]
            else:
                target_paragraph = target_text_frame.add_paragraph()
            
            # 复制段落级别的格式
            if hasattr(source_paragraph, 'alignment'):
                target_paragraph.alignment = source_paragraph.alignment
            if hasattr(source_paragraph, 'level'):
                target_paragraph.level = source_paragraph.level
            
            # 复制运行（文本片段）
            for source_run in source_paragraph.runs:
                target_run = target_paragraph.add_run()
                target_run.text = source_run.text
                
                # 复制字体格式
                if hasattr(source_run, 'font') and hasattr(target_run, 'font'):
                    source_font = source_run.font
                    target_font = target_run.font
                    
                    if source_font.name:
                        target_font.name = source_font.name
                    if source_font.size:
                        target_font.size = source_font.size
                    if source_font.bold is not None:
                        target_font.bold = source_font.bold
                    if source_font.italic is not None:
                        target_font.italic = source_font.italic
                    if source_font.underline is not None:
                        target_font.underline = source_font.underline
                    if hasattr(source_font, 'color') and source_font.color:
                        target_font.color.rgb = source_font.color.rgb
    
    except Exception as e:
        print(f"复制文本框架时出错: {str(e)}")


def copy_text_frame_formatting(source_text_frame, target_text_frame):
    """
    复制文本框架的格式设置（不包括文本内容）
    
    Args:
        source_text_frame: 源文本框架
        target_text_frame: 目标文本框架
    """
    try:
        # 复制文本框架的边距设置
        if hasattr(source_text_frame, 'margin_left') and hasattr(target_text_frame, 'margin_left'):
            target_text_frame.margin_left = source_text_frame.margin_left
        if hasattr(source_text_frame, 'margin_right') and hasattr(target_text_frame, 'margin_right'):
            target_text_frame.margin_right = source_text_frame.margin_right
        if hasattr(source_text_frame, 'margin_top') and hasattr(target_text_frame, 'margin_top'):
            target_text_frame.margin_top = source_text_frame.margin_top
        if hasattr(source_text_frame, 'margin_bottom') and hasattr(target_text_frame, 'margin_bottom'):
            target_text_frame.margin_bottom = source_text_frame.margin_bottom
        
        # 复制文本自动适应设置
        if hasattr(source_text_frame, 'auto_size') and hasattr(target_text_frame, 'auto_size'):
            target_text_frame.auto_size = source_text_frame.auto_size
        
        # 复制垂直对齐方式
        if hasattr(source_text_frame, 'vertical_anchor') and hasattr(target_text_frame, 'vertical_anchor'):
            target_text_frame.vertical_anchor = source_text_frame.vertical_anchor
        
        # 复制文字换行设置
        if hasattr(source_text_frame, 'word_wrap') and hasattr(target_text_frame, 'word_wrap'):
            target_text_frame.word_wrap = source_text_frame.word_wrap
            
    except Exception as e:
        print(f"复制文本框架格式时出错: {str(e)}")


def copy_shape_to_slide(source_shape, target_slide):
    """
    复制形状到目标幻灯片
    
    Args:
        source_shape: 源形状
        target_slide: 目标幻灯片
    """
    try:
        from pptx.shapes.autoshape import Shape
        
        # 根据形状类型进行复制
        if source_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # 复制文本框
            left = source_shape.left
            top = source_shape.top
            width = source_shape.width
            height = source_shape.height
            
            new_textbox = target_slide.shapes.add_textbox(left, top, width, height)
            
            # 复制文本内容和格式
            if hasattr(source_shape, 'text_frame'):
                copy_text_frame(source_shape.text_frame, new_textbox.text_frame)
        
        elif source_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # 复制自动形状
            # 这里可以扩展更多形状类型的复制逻辑
            pass
            
        # 可以继续添加其他形状类型的处理...
        
    except Exception as e:
        print(f"复制形状时出错: {str(e)}")


def duplicate_first_slide_multiple_times(file_path: str, copy_count: int = 2, output_file: str = None) -> Dict[str, any]:
    """
    专门用于复制第一张幻灯片多次的便捷函数
    
    Args:
        file_path (str): PowerPoint 文档的文件路径
        copy_count (int): 复制的数量，默认为2
        output_file (str): 输出文件路径，如果为None则生成默认输出文件名
        
    Returns:
        Dict[str, any]: 包含操作结果的字典
    """
    # 如果没有指定输出文件，生成默认文件名
    if output_file is None:
        base_name = os.path.splitext(file_path)[0]
        extension = os.path.splitext(file_path)[1]
        output_file = f"{base_name}_复制版{extension}"
    
    return duplicate_slide(file_path, 1, copy_count, output_file)


def print_duplication_summary(file_path: str, result: Dict[str, any]) -> None:
    """
    打印复制操作的摘要信息
    
    Args:
        file_path (str): PowerPoint 文档的文件路径
        result (Dict[str, any]): 复制操作的结果
    """
    print("📋 幻灯片复制操作摘要")
    print("=" * 60)
    print(f"源文件路径: {file_path}")
    
    if result['success']:
        print(f"✅ 复制成功!")
        print(f"📊 复制的幻灯片数量: {result['copied_slides']}")
        print(f"📈 操作后总幻灯片数: {result['total_slides_after']}")
        print(f"💾 输出文件: {result['output_file']}")
        
        # 如果输出文件不同于源文件，说明是另存为
        if result['output_file'] != file_path:
            print(f"✨ 已另存为新文件，原文件保持不变")
        else:
            print(f"⚠️  已覆盖原文件")
    else:
        print(f"❌ 复制失败: {result['error']}")


def create_slide_copies_with_new_filename(file_path: str, copy_count: int = 2, custom_suffix: str = None) -> Dict[str, any]:
    """
    创建幻灯片副本并保存为新文件（带自定义后缀）
    
    Args:
        file_path (str): PowerPoint 文档的文件路径
        copy_count (int): 复制的数量，默认为2
        custom_suffix (str): 自定义文件名后缀，如果为None则使用默认后缀
        
    Returns:
        Dict[str, any]: 包含操作结果的字典
    """
    # 生成输出文件名
    base_name = os.path.splitext(file_path)[0]
    extension = os.path.splitext(file_path)[1]
    
    if custom_suffix:
        output_file = f"{base_name}_{custom_suffix}{extension}"
    else:
        output_file = f"{base_name}_复制版{extension}"
    
    return duplicate_first_slide_multiple_times(file_path, copy_count, output_file)


def replace_slides_with_word_sections(word_file_path: str, ppt_file_path: str, output_file: str = None) -> Dict[str, any]:
    """
    將Word文檔的分段內容替換到PowerPoint幻燈片中，每一段作為一頁投影片
    增強版本：更好的格式保留和背景處理
    
    Args:
        word_file_path (str): Word文檔的文件路徑
        ppt_file_path (str): PowerPoint文檔的文件路徑
        output_file (str): 輸出文件路徑，如果為None則生成默認輸出文件名
        
    Returns:
        Dict[str, any]: 包含操作結果的字典
    """
    result = {
        'success': False,
        'error': None,
        'total_sections': 0,
        'total_slides_created': 0,
        'output_file': output_file,
        'skipped_sections': [],
        'format_issues': []
    }
    
    try:
        # 1. 解析Word文檔的分段內容
        print("🔍 正在解析Word文檔的分段內容...")
        sections_result = parse_numbered_sections(word_file_path)
        if not sections_result['success']:
            result['error'] = f"解析Word文檔失敗: {sections_result['error']}"
            return result
        
        sections = sections_result['sections']
        result['total_sections'] = len(sections)
        print(f"✅ 找到 {len(sections)} 個分段")
        
        # 2. 檢查PowerPoint文檔是否存在
        if not os.path.exists(ppt_file_path):
            result['error'] = f"PowerPoint文件不存在: {ppt_file_path}"
            return result
        
        # 3. 讀取PowerPoint演示文稿
        print("📖 正在讀取PowerPoint模板...")
        prs = Presentation(ppt_file_path)
        
        # 獲取第一張幻燈片作為模板
        if len(prs.slides) == 0:
            result['error'] = "PowerPoint文件中沒有幻燈片可作為模板"
            return result
        
        template_slide = prs.slides[0]
        template_layout = template_slide.slide_layout
        
        # 分析模板幻燈片的格式
        print("🔍 正在分析模板幻燈片格式...")
        template_analysis = analyze_template_slide(template_slide)
        print(f"   📊 模板分析: {template_analysis['summary']}")
        
        # 4. 清除現有的幻燈片，保留第一張作為模板
        print("🧹 正在清理現有幻燈片...")
        # 從後往前刪除，避免索引問題
        for i in range(len(prs.slides) - 1, 0, -1):
            slide_to_remove = prs.slides[i]
            rId = prs.slides._slides[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._slides[i]
        
        # 5. 為每個分段創建一張幻燈片
        print("📝 正在為每個分段創建幻燈片...")
        slides_created = 0
        
        for i, section in enumerate(sections):
            try:
                print(f"   處理段落 {section['number']}: {section['title'][:50]}...")
                
                if i == 0:
                    # 第一個分段使用現有的第一張幻燈片
                    target_slide = template_slide
                    print(f"     使用第一張幻燈片作為模板")
                else:
                    # 其他分段創建新幻燈片
                    print(f"     創建新幻燈片...")
                    target_slide = prs.slides.add_slide(template_layout)
                    
                    # 使用增強的複製方法
                    print(f"     複製模板格式...")
                    copy_slide_with_enhanced_formatting(template_slide, target_slide, template_analysis)
                
                # 替換幻燈片內容
                print(f"     替換內容...")
                replace_slide_content_with_section(target_slide, section, template_slide)
                slides_created += 1
                print(f"     ✅ 段落 {section['number']} 處理完成")
                
            except Exception as e:
                print(f"⚠️  處理段落 {section['number']} 時出錯: {str(e)}")
                result['skipped_sections'].append({
                    'number': section['number'],
                    'title': section['title'],
                    'error': str(e)
                })
                continue
        
        result['total_slides_created'] = slides_created
        
        # 6. 生成輸出文件名（如果未指定）
        if output_file is None:
            base_name = os.path.splitext(ppt_file_path)[0]
            extension = os.path.splitext(ppt_file_path)[1]
            output_file = f"{base_name}_分段版{extension}"
        
        result['output_file'] = output_file
        
        # 7. 保存文件
        print(f"💾 正在保存到 {output_file}...")
        prs.save(output_file)
        
        result['success'] = True
        print(f"✅ 成功！創建了 {slides_created} 張幻燈片")
        
        # 8. 輸出格式保留情況摘要
        if result['format_issues']:
            print(f"⚠️  發現 {len(result['format_issues'])} 個格式問題")
            for issue in result['format_issues'][:3]:  # 只顯示前3個
                print(f"     - {issue}")
        else:
            print(f"✅ 格式保留良好")
        
    except Exception as e:
        result['error'] = f"處理過程中發生錯誤: {str(e)}"
    
    return result


def analyze_template_slide(slide):
    """
    分析模板幻燈片的格式特徵
    
    Args:
        slide: 模板幻燈片
        
    Returns:
        Dict: 分析結果
    """
    analysis = {
        'text_shapes': [],
        'background_shapes': [],
        'has_background_image': False,
        'layout_name': '',
        'master_name': '',
        'summary': ''
    }
    
    try:
        # 分析佈局信息
        if hasattr(slide, 'slide_layout'):
            analysis['layout_name'] = slide.slide_layout.name
            if hasattr(slide.slide_layout, 'slide_master'):
                analysis['master_name'] = getattr(slide.slide_layout.slide_master, 'name', 'Unknown')
        
        # 分析形狀
        text_count = 0
        image_count = 0
        other_count = 0
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and hasattr(shape, 'text'):
                text_count += 1
                analysis['text_shapes'].append({
                    'type': 'text',
                    'has_text': bool(shape.text.strip()),
                    'width': shape.width,
                    'height': shape.height,
                    'left': shape.left,
                    'top': shape.top
                })
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                # 檢查是否可能是背景圖片
                slide_width = slide.slide_layout.width if hasattr(slide.slide_layout, 'width') else 0
                slide_height = slide.slide_layout.height if hasattr(slide.slide_layout, 'height') else 0
                
                is_background = False
                if (slide_width > 0 and slide_height > 0 and
                    shape.width >= slide_width * 0.8 and 
                    shape.height >= slide_height * 0.8):
                    is_background = True
                    analysis['has_background_image'] = True
                
                analysis['background_shapes'].append({
                    'type': 'picture',
                    'is_background': is_background,
                    'width': shape.width,
                    'height': shape.height
                })
            else:
                other_count += 1
        
        # 生成摘要
        summary_parts = []
        if text_count > 0:
            summary_parts.append(f"{text_count}個文本框")
        if image_count > 0:
            summary_parts.append(f"{image_count}個圖片")
        if other_count > 0:
            summary_parts.append(f"{other_count}個其他形狀")
        if analysis['has_background_image']:
            summary_parts.append("有背景圖片")
        
        analysis['summary'] = ', '.join(summary_parts) if summary_parts else '無特殊元素'
        
    except Exception as e:
        analysis['summary'] = f"分析失敗: {str(e)}"
    
    return analysis


def copy_slide_with_enhanced_formatting(source_slide, target_slide, template_analysis):
    """
    使用增強的格式複製方法
    
    Args:
        source_slide: 源幻燈片
        target_slide: 目標幻燈片
        template_analysis: 模板分析結果
    """
    try:
        print(f"       🎨 使用增強格式複製...")
        
        # 1. 複製背景（包括背景圖片）
        copy_slide_background(source_slide, target_slide)
        
        # 2. 按順序複製形狀（保持Z順序）
        for i, shape in enumerate(source_slide.shapes):
            try:
                shape_info = f"形狀{i+1}"
                if hasattr(shape, 'name'):
                    shape_info += f"({shape.name})"
                    
                print(f"         複製{shape_info}...")
                
                if shape.is_placeholder:
                    # 處理占位符
                    copy_placeholder_content(shape, target_slide)
                else:
                    # 處理非占位符形狀
                    copy_non_placeholder_shape(shape, target_slide)
                    
            except Exception as shape_error:
                print(f"         ⚠️  複製{shape_info}失敗: {str(shape_error)}")
                continue
        
        print(f"       ✅ 增強格式複製完成")
        
    except Exception as e:
        print(f"       ❌ 增強格式複製失敗: {str(e)}")
        # 如果增強複製失敗，回退到基本複製
        try:
            copy_slide_completely(source_slide, target_slide)
            print(f"       ✅ 回退到基本複製完成")
        except Exception as fallback_error:
            print(f"       ❌ 基本複製也失敗: {str(fallback_error)}")


def replace_slide_content_with_section(slide, section, template_slide=None):
    """
    將分段內容替換到指定幻燈片中，保持原有的格式
    
    Args:
        slide: 目標幻燈片
        section: 分段數據（包含number, title, content等）
        template_slide: 模板幻燈片（用於獲取格式信息）
    """
    try:
        # 查找幻灯片中的文本框
        text_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and hasattr(shape, 'text'):
                text_shapes.append(shape)
        
        # 如果没有找到文本框，从模板复制或创建一个新的文本框
        if not text_shapes:
            print(f"   ⚠️  幻燈片中沒有找到文本框...")
            
            # 如果有模板幻灯片，尝试从模板复制文本框
            if template_slide:
                template_text_shapes = []
                for shape in template_slide.shapes:
                    if hasattr(shape, 'text_frame') and hasattr(shape, 'text'):
                        template_text_shapes.append(shape)
                
                if template_text_shapes:
                    print(f"   📋 正在從模板複製文本框格式...")
                    # 复制第一个文本框的位置和大小
                    template_shape = template_text_shapes[0]
                    new_textbox = slide.shapes.add_textbox(
                        template_shape.left, 
                        template_shape.top, 
                        template_shape.width, 
                        template_shape.height
                    )
                    
                    # 复制文本框的格式（填充、线条等）
                    copy_shape_formatting(template_shape, new_textbox)
                    
                    # 复制文本框架的格式
                    if hasattr(template_shape, 'text_frame') and hasattr(new_textbox, 'text_frame'):
                        copy_text_frame_formatting(template_shape.text_frame, new_textbox.text_frame)
                    
                    text_shapes.append(new_textbox)
                    print(f"   ✅ 成功從模板複製文本框格式")
                else:
                    # 模板也没有文本框，创建默认的
                    from pptx.util import Inches
                    new_textbox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1), Inches(9), Inches(6.5)
                    )
                    text_shapes.append(new_textbox)
                    print(f"   ✅ 成功創建默認文本框")
            else:
                # 没有模板，创建默认的文本框
                from pptx.util import Inches
                new_textbox = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1), Inches(9), Inches(6.5)
                )
                text_shapes.append(new_textbox)
                print(f"   ✅ 成功創建默認文本框")
        
        # 準備要替換的內容
        # 使用段落編號作為標題（如果不是前言）
        if section['number'] == 0:
            title_text = section['title']  # 前言
        else:
            title_text = f"{section['number']}. {section['title']}"
        
        # 內容文本（去除第一行，因為第一行通常是標題）
        content_lines = section['content']
        if len(content_lines) > 1:
            content_text = '\n'.join(content_lines[1:])  # 跳過第一行標題
        else:
            content_text = section['text_only']  # 如果只有一行，使用純文本
        
        # 如果內容為空，使用標題作為內容
        if not content_text.strip():
            content_text = title_text
        
        # 替換第一個文本框的內容，保持原有格式
        main_text_shape = text_shapes[0]
        
        # 詳細保存原有的文本格式信息
        original_formats = extract_detailed_text_formatting(main_text_shape)
        
        # 清空現有內容
        main_text_shape.text_frame.clear()
        
        # 設置新內容並應用原有格式，避免不必要的空行
        if content_text != title_text and content_text.strip():
            # 有標題和內容，只用一個換行符分隔
            full_text = f"{title_text}\n{content_text}"
        else:
            # 只有標題或內容相同
            full_text = title_text
        
        # 分段設置文本，以便保持不同層級的格式，同時過濾空行
        lines = [line for line in full_text.split('\n') if line.strip()]  # 過濾掉空行
        
        for i, line in enumerate(lines):
            if i == 0:
                # 第一段，使用現有的段落
                paragraph = main_text_shape.text_frame.paragraphs[0]
            else:
                # 其他段落，添加新段落
                paragraph = main_text_shape.text_frame.add_paragraph()
            
            # 添加文本運行並應用格式
            run = paragraph.add_run()
            run.text = line
            
            # 應用格式：標題行使用標題格式，內容行使用內容格式
            if i == 0:  # 標題行
                apply_text_formatting(run, paragraph, original_formats.get('title', original_formats.get('default')))
            else:  # 內容行
                apply_text_formatting(run, paragraph, original_formats.get('content', original_formats.get('default')))
        
        # 如果有多個文本框，清空其他文本框
        for i in range(1, len(text_shapes)):
            text_shapes[i].text = ""
        
        print(f"   ✅ 成功替換段落 {section['number']} 的內容，保持格式")
        
    except Exception as e:
        print(f"   ❌ 替換段落 {section['number']} 內容時出錯: {str(e)}")
        raise


def extract_detailed_text_formatting(text_shape):
    """
    提取文本框中的詳細格式信息
    
    Args:
        text_shape: 文本形狀
        
    Returns:
        Dict: 包含各種格式信息的字典
    """
    formats = {
        'default': None,
        'title': None,
        'content': None
    }
    
    try:
        if not hasattr(text_shape, 'text_frame') or not text_shape.text_frame.paragraphs:
            return formats
        
        # 分析不同段落的格式
        paragraphs = text_shape.text_frame.paragraphs
        
        for i, paragraph in enumerate(paragraphs):
            if not paragraph.runs:
                continue
                
            first_run = paragraph.runs[0]
            format_info = {
                'font_name': getattr(first_run.font, 'name', None),
                'font_size': getattr(first_run.font, 'size', None),
                'font_bold': getattr(first_run.font, 'bold', None),
                'font_italic': getattr(first_run.font, 'italic', None),
                'font_underline': getattr(first_run.font, 'underline', None),
                'font_color': None,
                'alignment': getattr(paragraph, 'alignment', None),
                'level': getattr(paragraph, 'level', None)
            }
            
            # 嘗試獲取字體顏色
            try:
                if hasattr(first_run.font, 'color') and hasattr(first_run.font.color, 'rgb'):
                    format_info['font_color'] = first_run.font.color.rgb
            except:
                pass
            
            # 根據位置和特徵判斷格式類型
            if i == 0:
                # 第一段通常是標題
                formats['title'] = format_info
                if formats['default'] is None:
                    formats['default'] = format_info
            else:
                # 其他段落是內容
                if formats['content'] is None:
                    formats['content'] = format_info
            
            # 設置默認格式
            if formats['default'] is None:
                formats['default'] = format_info
        
        # 如果沒有提取到格式，創建基本默認格式
        if formats['default'] is None:
            from pptx.util import Pt
            formats['default'] = {
                'font_name': 'Arial',
                'font_size': Pt(24),
                'font_bold': False,
                'font_italic': False,
                'font_underline': None,
                'font_color': None,
                'alignment': None,
                'level': 0
            }
        
        # 確保標題和內容格式存在
        if formats['title'] is None:
            formats['title'] = formats['default'].copy()
            formats['title']['font_bold'] = True  # 標題通常是粗體
        
        if formats['content'] is None:
            formats['content'] = formats['default'].copy()
            
    except Exception as e:
        print(f"     ⚠️  提取文本格式時出錯: {str(e)}")
        # 返回基本格式
        from pptx.util import Pt
        basic_format = {
            'font_name': 'Arial',
            'font_size': Pt(24),
            'font_bold': False,
            'font_italic': False,
            'font_underline': None,
            'font_color': None,
            'alignment': None,
            'level': 0
        }
        formats = {
            'default': basic_format,
            'title': basic_format.copy(),
            'content': basic_format.copy()
        }
        formats['title']['font_bold'] = True
    
    return formats


def apply_text_formatting(run, paragraph, format_info):
    """
    將格式信息應用到文本運行和段落
    
    Args:
        run: 文本運行對象
        paragraph: 段落對象
        format_info: 格式信息字典
    """
    if not format_info:
        return
        
    try:
        # 應用字體格式
        if format_info.get('font_name'):
            run.font.name = format_info['font_name']
        
        if format_info.get('font_size'):
            run.font.size = format_info['font_size']
        
        if format_info.get('font_bold') is not None:
            run.font.bold = format_info['font_bold']
        
        if format_info.get('font_italic') is not None:
            run.font.italic = format_info['font_italic']
        
        if format_info.get('font_underline') is not None:
            run.font.underline = format_info['font_underline']
        
        if format_info.get('font_color'):
            run.font.color.rgb = format_info['font_color']
        
        # 應用段落格式
        if format_info.get('alignment') is not None:
            paragraph.alignment = format_info['alignment']
        
        if format_info.get('level') is not None:
            paragraph.level = format_info['level']
            
    except Exception as e:
        print(f"     ⚠️  應用文本格式時出錯: {str(e)}")


def print_replacement_summary(word_file: str, ppt_file: str, result: Dict[str, any]) -> None:
    """
    打印替換操作的摘要信息
    
    Args:
        word_file (str): Word文檔路徑
        ppt_file (str): PowerPoint文檔路徑  
        result (Dict[str, any]): 替換操作的結果
    """
    print("\n📋 Word轉PowerPoint操作摘要")
    print("=" * 70)
    print(f"Word文檔: {word_file}")
    print(f"PowerPoint模板: {ppt_file}")
    
    if result['success']:
        print(f"✅ 轉換成功!")
        print(f"📊 Word分段數量: {result['total_sections']}")
        print(f"📈 創建的幻燈片數: {result['total_slides_created']}")
        print(f"💾 輸出文件: {result['output_file']}")
        
        if result['skipped_sections']:
            print(f"⚠️  跳過的段落數: {len(result['skipped_sections'])}")
            print("跳過的段落:")
            for skipped in result['skipped_sections']:
                print(f"   段落 {skipped['number']}: {skipped['title'][:50]}... (錯誤: {skipped['error']})")
        
        print(f"\n🎯 請檢查輸出文件: {result['output_file']}")
    else:
        print(f"❌ 轉換失敗: {result['error']}")


if __name__ == "__main__":
    # 文件路径
    word_file = "證道資料.docx"
    ppt_file = "證道資料.pptx"
    
    print("=" * 80)
    print("测试Word文档读取功能")
    print("=" * 80)
    
    if os.path.exists(word_file):
        # 打印Word文档摘要
        print_word_summary(word_file)
        
        print("\n" + "=" * 60)
        print("Word文档段落解析")
        print("=" * 60)
        
        # 解析并打印段落摘要
        print_sections_summary(word_file)
    else:
        print(f"❌ Word文件不存在: {word_file}")
    
    print("\n" + "=" * 80)
    print("测试PowerPoint文档读取功能")
    print("=" * 80)
    
    if os.path.exists(ppt_file):
        # 打印PowerPoint文档摘要
        print_powerpoint_summary(ppt_file)
        
        print("\n" + "=" * 60)
        print("PowerPoint详细内容预览")
        print("=" * 60)
        
        # 查看前几页幻灯片的详细内容（根据实际幻灯片数量决定）
        ppt = Presentation(ppt_file)
        total_slides = len(ppt.slides)
        slides_to_view = min(3, total_slides)  # 最多查看3页，但不超过实际页数
        
        for slide_num in range(1, slides_to_view + 1):
            print_slide_detail(ppt_file, slide_num)
            print()
        
        # 测试Word分段内容替换到PowerPoint功能
        print("\n" + "=" * 60)
        print("测试Word分段内容替换到PowerPoint功能")
        print("=" * 60)
        
        if os.path.exists(word_file):
            print("🔄 正在将Word分段内容替换到PowerPoint中...")
            replacement_result = replace_slides_with_word_sections(word_file, ppt_file)
            print_replacement_summary(word_file, ppt_file, replacement_result)
            
            if replacement_result['success']:
                print(f"\n📋 轉換後的PowerPoint摘要:")
                print_powerpoint_summary(replacement_result['output_file'])
        else:
            print(f"❌ Word文件不存在，無法進行轉換: {word_file}")
    else:
        print(f"❌ PowerPoint文件不存在: {ppt_file}")
    
    print("=" * 80)
