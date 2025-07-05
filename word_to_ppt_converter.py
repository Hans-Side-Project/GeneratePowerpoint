"""
Word转PowerPoint转换器
优化版本：专注于核心功能，移除冗余代码和日志
"""

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from typing import List, Dict, Optional
import os
import re
import io


def extract_word_formatting(paragraph):
    """
    从Word段落中提取详细的格式信息，包括每个文本片段的颜色
    """
    formatting_data = []
    
    for run in paragraph.runs:
        run_format = {
            'text': run.text,
            'font_name': None,
            'font_size': None,
            'font_bold': None,
            'font_italic': None,
            'font_underline': None,
            'font_color': None
        }
        
        if hasattr(run, 'font'):
            font = run.font
            run_format['font_name'] = font.name
            run_format['font_size'] = font.size
            run_format['font_bold'] = font.bold
            run_format['font_italic'] = font.italic
            run_format['font_underline'] = font.underline
            
            # 提取颜色信息
            try:
                if hasattr(font, 'color') and font.color:
                    if hasattr(font.color, 'rgb') and font.color.rgb:
                        run_format['font_color'] = font.color.rgb
                    elif hasattr(font.color, 'theme_color') and font.color.theme_color:
                        # 处理主题颜色
                        run_format['font_color'] = 'theme_color'
            except Exception:
                pass
        
        formatting_data.append(run_format)
    
    return formatting_data


def parse_word_sections(file_path: str) -> Dict[str, any]:
    """
    解析Word文档的编号段落，保留详细的格式信息
    
    Args:
        file_path (str): Word文档路径
        
    Returns:
        Dict: 包含分段内容和格式信息的字典
    """
    result = {
        'sections': [],
        'total_sections': 0,
        'success': False,
        'error': None
    }
    
    try:
        if not os.path.exists(file_path):
            result['error'] = f"文件不存在: {file_path}"
            return result
            
        doc = Document(file_path)
        sections = []
        current_section = None
        number_pattern = re.compile(r'^(\d+)\.\s*(.*)')
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            
            # 提取格式信息
            paragraph_formatting = extract_word_formatting(paragraph)
            
            match = number_pattern.match(text)
            if match:
                if current_section is not None:
                    sections.append(current_section)
                
                section_number = int(match.group(1))
                section_title = match.group(2) if match.group(2) else ""
                
                current_section = {
                    'number': section_number,
                    'title': section_title,
                    'content': [text],
                    'text_only': section_title,
                    'formatting': [paragraph_formatting]
                }
            else:
                if current_section is not None:
                    current_section['content'].append(text)
                    current_section['formatting'].append(paragraph_formatting)
                    if current_section['text_only']:
                        current_section['text_only'] += '\n' + text
                    else:
                        current_section['text_only'] = text
                else:
                    current_section = {
                        'number': 0,
                        'title': '前言',
                        'content': [text],
                        'text_only': text,
                        'formatting': [paragraph_formatting]
                    }
        
        if current_section is not None:
            sections.append(current_section)
        
        result['sections'] = sections
        result['total_sections'] = len(sections)
        result['success'] = True
        
    except Exception as e:
        result['error'] = f"解析Word文档失败: {str(e)}"
    
    return result


def copy_shape_formatting(source_shape, target_shape):
    """
    复制形状格式，优化版本
    """
    try:
        # 基本几何属性
        target_shape.left = source_shape.left
        target_shape.top = source_shape.top
        target_shape.width = source_shape.width
        target_shape.height = source_shape.height
        
        # 填充格式
        if hasattr(source_shape, 'fill') and hasattr(target_shape, 'fill'):
            try:
                if hasattr(source_shape.fill, 'type'):
                    fill_type = source_shape.fill.type
                    if fill_type == MSO_FILL_TYPE.SOLID:
                        target_shape.fill.solid()
                        if hasattr(source_shape.fill, 'fore_color') and hasattr(source_shape.fill.fore_color, 'rgb'):
                            target_shape.fill.fore_color.rgb = source_shape.fill.fore_color.rgb
                    else:
                        target_shape.fill.background()
                else:
                    target_shape.fill.background()
            except:
                target_shape.fill.background()
        
        # 线条格式
        if hasattr(source_shape, 'line') and hasattr(target_shape, 'line'):
            try:
                line_width_pt = 0
                if (hasattr(source_shape.line, 'width') and 
                    source_shape.line.width is not None and 
                    hasattr(source_shape.line.width, 'pt')):
                    line_width_pt = source_shape.line.width.pt
                
                if line_width_pt > 0:
                    target_shape.line.width = source_shape.line.width
                    if hasattr(source_shape.line, 'color') and hasattr(target_shape.line, 'color'):
                        if hasattr(source_shape.line.color, 'rgb'):
                            target_shape.line.color.rgb = source_shape.line.color.rgb
                else:
                    from pptx.util import Pt
                    target_shape.line.width = Pt(0)
            except:
                pass
                
    except Exception:
        pass


def copy_text_frame(source_text_frame, target_text_frame):
    """
    复制文本框架内容和格式
    """
    try:
        target_text_frame.clear()
        
        for source_paragraph in source_text_frame.paragraphs:
            if len(target_text_frame.paragraphs) == 1 and not target_text_frame.paragraphs[0].text:
                target_paragraph = target_text_frame.paragraphs[0]
            else:
                target_paragraph = target_text_frame.add_paragraph()
            
            if hasattr(source_paragraph, 'alignment'):
                target_paragraph.alignment = source_paragraph.alignment
            if hasattr(source_paragraph, 'level'):
                target_paragraph.level = source_paragraph.level
            
            for source_run in source_paragraph.runs:
                target_run = target_paragraph.add_run()
                target_run.text = source_run.text
                
                if hasattr(source_run, 'font') and hasattr(target_run, 'font'):
                    source_font = source_run.font
                    target_font = target_run.font
                    
                    if source_font.name:
                        target_font.name = source_font.name
                    if source_font.size:
                        target_font.size = source_font.size
                    if source_font.bold is not None:
                        target_font.bold = source_font.bold
                    if source_font.italic is not None:
                        target_font.italic = source_font.italic
                    if source_font.underline is not None:
                        target_font.underline = source_font.underline
                    if hasattr(source_font, 'color') and source_font.color:
                        target_font.color.rgb = source_font.color.rgb
    except Exception:
        pass


def copy_slide_content(source_slide, target_slide):
    """
    复制幻灯片内容，优化版本
    """
    try:
        for shape in source_slide.shapes:
            try:
                if shape.is_placeholder:
                    # 处理占位符
                    for target_shape in target_slide.shapes:
                        if (target_shape.is_placeholder and 
                            hasattr(shape, 'placeholder_format') and
                            hasattr(target_shape, 'placeholder_format') and
                            target_shape.placeholder_format.idx == shape.placeholder_format.idx):
                            
                            if hasattr(shape, 'text') and hasattr(target_shape, 'text'):
                                target_shape.text = shape.text
                            
                            if hasattr(shape, 'text_frame') and hasattr(target_shape, 'text_frame'):
                                copy_text_frame(shape.text_frame, target_shape.text_frame)
                            break
                else:
                    # 处理非占位符形状
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    
                    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        new_textbox = target_slide.shapes.add_textbox(left, top, width, height)
                        if hasattr(shape, 'text_frame'):
                            copy_text_frame(shape.text_frame, new_textbox.text_frame)
                        copy_shape_formatting(shape, new_textbox)
                        
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_blob = shape.image.blob
                            image_stream = io.BytesIO(image_blob)
                            new_picture = target_slide.shapes.add_picture(image_stream, left, top, width, height)
                            copy_shape_formatting(shape, new_picture)
                        except:
                            pass
                            
                    elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        try:
                            if hasattr(shape, 'auto_shape_type'):
                                new_shape = target_slide.shapes.add_shape(
                                    shape.auto_shape_type, left, top, width, height
                                )
                                if hasattr(shape, 'text_frame'):
                                    copy_text_frame(shape.text_frame, new_shape.text_frame)
                                copy_shape_formatting(shape, new_shape)
                        except:
                            pass
            except:
                continue
    except Exception:
        pass
    except Exception:
        pass


def extract_text_formatting(text_shape):
    """
    提取文本格式，简化版本
    """
    formats = {'default': None, 'title': None, 'content': None}
    
    try:
        if not hasattr(text_shape, 'text_frame') or not text_shape.text_frame.paragraphs:
            return formats
        
        paragraphs = text_shape.text_frame.paragraphs
        
        for i, paragraph in enumerate(paragraphs):
            if not paragraph.runs:
                continue
                
            first_run = paragraph.runs[0]
            format_info = {
                'font_name': getattr(first_run.font, 'name', None),
                'font_size': getattr(first_run.font, 'size', None),
                'font_bold': getattr(first_run.font, 'bold', None),
                'font_italic': getattr(first_run.font, 'italic', None),
                'font_underline': getattr(first_run.font, 'underline', None),
                'font_color': None,
                'alignment': getattr(paragraph, 'alignment', None),
                'level': getattr(paragraph, 'level', None)
            }
            
            try:
                if hasattr(first_run.font, 'color') and hasattr(first_run.font.color, 'rgb'):
                    format_info['font_color'] = first_run.font.color.rgb
            except:
                pass
            
            if i == 0:
                formats['title'] = format_info
                if formats['default'] is None:
                    formats['default'] = format_info
            else:
                if formats['content'] is None:
                    formats['content'] = format_info
            
            if formats['default'] is None:
                formats['default'] = format_info
        
        # 确保有基本格式
        if formats['default'] is None:
            from pptx.util import Pt
            formats['default'] = {
                'font_name': 'Arial',
                'font_size': Pt(24),
                'font_bold': False,
                'font_italic': False,
                'font_underline': None,
                'font_color': None,
                'alignment': None,
                'level': 0
            }
        
        if formats['title'] is None:
            formats['title'] = formats['default'].copy()
            formats['title']['font_bold'] = True
        
        if formats['content'] is None:
            formats['content'] = formats['default'].copy()
            
    except Exception:
        from pptx.util import Pt
        basic_format = {
            'font_name': 'Arial',
            'font_size': Pt(24),
            'font_bold': False,
            'font_italic': False,
            'font_underline': None,
            'font_color': None,
            'alignment': None,
            'level': 0
        }
        formats = {
            'default': basic_format,
            'title': basic_format.copy(),
            'content': basic_format.copy()
        }
        formats['title']['font_bold'] = True
    
    return formats


def apply_text_formatting(run, paragraph, format_info):
    """
    应用文本格式，简化版本
    """
    if not format_info:
        return
        
    try:
        if format_info.get('font_name'):
            run.font.name = format_info['font_name']
        if format_info.get('font_size'):
            run.font.size = format_info['font_size']
        if format_info.get('font_bold') is not None:
            run.font.bold = format_info['font_bold']
        if format_info.get('font_italic') is not None:
            run.font.italic = format_info['font_italic']
        if format_info.get('font_underline') is not None:
            run.font.underline = format_info['font_underline']
        if format_info.get('font_color'):
            run.font.color.rgb = format_info['font_color']
        if format_info.get('alignment') is not None:
            paragraph.alignment = format_info['alignment']
        if format_info.get('level') is not None:
            paragraph.level = format_info['level']
    except Exception:
        pass


def apply_word_formatting_to_run(run, format_info):
    """
    将Word格式信息应用到PowerPoint文本运行
    """
    try:
        if format_info.get('font_name'):
            run.font.name = format_info['font_name']
        if format_info.get('font_size'):
            run.font.size = format_info['font_size']
        if format_info.get('font_bold') is not None:
            run.font.bold = format_info['font_bold']
        if format_info.get('font_italic') is not None:
            run.font.italic = format_info['font_italic']
        if format_info.get('font_underline') is not None:
            run.font.underline = format_info['font_underline']
        
        # 应用颜色
        if format_info.get('font_color'):
            try:
                if format_info['font_color'] != 'theme_color':
                    run.font.color.rgb = format_info['font_color']
            except Exception:
                pass
                
    except Exception:
        pass


def replace_slide_content_with_formatting(slide, section, template_slide=None):
    """
    替换幻灯片内容并保留Word文档的格式
    """
    try:
        # 查找文本框
        text_shapes = [shape for shape in slide.shapes 
                      if hasattr(shape, 'text_frame') and hasattr(shape, 'text')]
        
        # 如果没有文本框，从模板创建
        if not text_shapes and template_slide:
            template_text_shapes = [shape for shape in template_slide.shapes 
                                  if hasattr(shape, 'text_frame') and hasattr(shape, 'text')]
            
            if template_text_shapes:
                template_shape = template_text_shapes[0]
                new_textbox = slide.shapes.add_textbox(
                    template_shape.left, template_shape.top, 
                    template_shape.width, template_shape.height
                )
                copy_shape_formatting(template_shape, new_textbox)
                text_shapes.append(new_textbox)
            else:
                from pptx.util import Inches
                new_textbox = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1), Inches(9), Inches(6.5)
                )
                text_shapes.append(new_textbox)
        
        if not text_shapes:
            return
        
        # 主文本框
        main_text_shape = text_shapes[0]
        main_text_shape.text_frame.clear()
        
        # 处理每个段落的格式
        for content_idx, (content_line, formatting_data) in enumerate(zip(section['content'], section.get('formatting', []))):
            
            # 创建段落
            if content_idx == 0:
                paragraph = main_text_shape.text_frame.paragraphs[0]
            else:
                paragraph = main_text_shape.text_frame.add_paragraph()
            
            # 为每个格式化的文本片段创建运行
            for format_info in formatting_data:
                if format_info['text'].strip():  # 只处理非空文本
                    run = paragraph.add_run()
                    run.text = format_info['text']
                    apply_word_formatting_to_run(run, format_info)
        
        # 清空其他文本框
        for i in range(1, len(text_shapes)):
            text_shapes[i].text = ""
            
    except Exception as e:
        # 如果格式化失败，回退到基本文本设置
        try:
            main_text_shape = text_shapes[0] if text_shapes else None
            if main_text_shape:
                # 准备基本内容
                if section['number'] == 0:
                    title_text = section['title']
                else:
                    title_text = f"{section['number']}. {section['title']}"
                
                content_lines = section['content']
                if len(content_lines) > 1:
                    content_text = '\n'.join(content_lines[1:])
                else:
                    content_text = section['text_only']
                
                if not content_text.strip():
                    content_text = title_text
                
                # 设置基本文本
                if content_text != title_text and content_text.strip():
                    full_text = f"{title_text}\n{content_text}"
                else:
                    full_text = title_text
                
                # 过滤空行
                lines = [line for line in full_text.split('\n') if line.strip()]
                main_text_shape.text = '\n'.join(lines)
        except:
            pass


def convert_word_to_ppt(word_file_path: str, ppt_file_path: str, output_file: str = None) -> Dict[str, any]:
    """
    将Word文档转换为PowerPoint，优化版本
    
    Args:
        word_file_path (str): Word文档路径
        ppt_file_path (str): PowerPoint模板路径
        output_file (str): 输出文件路径
        
    Returns:
        Dict: 转换结果
    """
    result = {
        'success': False,
        'error': None,
        'total_sections': 0,
        'total_slides_created': 0,
        'output_file': output_file
    }
    
    try:
        # 解析Word文档
        sections_result = parse_word_sections(word_file_path)
        if not sections_result['success']:
            result['error'] = f"解析Word文档失败: {sections_result['error']}"
            return result
        
        sections = sections_result['sections']
        result['total_sections'] = len(sections)
        
        if not os.path.exists(ppt_file_path):
            result['error'] = f"PowerPoint模板不存在: {ppt_file_path}"
            return result
        
        # 读取PowerPoint
        prs = Presentation(ppt_file_path)
        
        if len(prs.slides) == 0:
            result['error'] = "PowerPoint模板中没有幻灯片"
            return result
        
        template_slide = prs.slides[0]
        template_layout = template_slide.slide_layout
        
        # 清除现有幻灯片（保留第一张）
        for i in range(len(prs.slides) - 1, 0, -1):
            slide_to_remove = prs.slides[i]
            rId = prs.slides._slides[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._slides[i]
        
        # 为每个段落创建幻灯片
        slides_created = 0
        
        for i, section in enumerate(sections):
            try:
                if i == 0:
                    target_slide = template_slide
                else:
                    target_slide = prs.slides.add_slide(template_layout)
                    copy_slide_content(template_slide, target_slide)
                
                replace_slide_content_with_formatting(target_slide, section, template_slide)
                slides_created += 1
                
            except Exception:
                continue
        
        result['total_slides_created'] = slides_created
        
        # 生成输出文件名
        if output_file is None:
            base_name = os.path.splitext(ppt_file_path)[0]
            extension = os.path.splitext(ppt_file_path)[1]
            output_file = f"{base_name}_分段版{extension}"
        
        result['output_file'] = output_file
        
        # 保存文件
        prs.save(output_file)
        result['success'] = True
        
    except Exception as e:
        result['error'] = f"转换过程出错: {str(e)}"
    
    return result


def main():
    """
    主函数：执行Word到PowerPoint的转换
    """
    word_file = "證道資料.docx"
    ppt_file = "證道資料.pptx"
    
    print("🔄 开始转换Word文档到PowerPoint...")
    
    if not os.path.exists(word_file):
        print(f"❌ Word文件不存在: {word_file}")
        return
    
    if not os.path.exists(ppt_file):
        print(f"❌ PowerPoint模板不存在: {ppt_file}")
        return
    
    result = convert_word_to_ppt(word_file, ppt_file)
    
    if result['success']:
        print(f"✅ 转换成功!")
        print(f"📊 处理段落数: {result['total_sections']}")
        print(f"📈 创建幻灯片数: {result['total_slides_created']}")
        print(f"💾 输出文件: {result['output_file']}")
    else:
        print(f"❌ 转换失败: {result['error']}")


if __name__ == "__main__":
    main()
