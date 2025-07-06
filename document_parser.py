"""
文件解析模組 - 負責解析 Word 和 PowerPoint 文件
提供統一的文件讀取、解析和結構化處理介面
"""

from typing import Dict, List, Any, Optional
from docx import Document
from pptx import Presentation
import os
import re
import logging
from format_handler import FormatHandler


class DocumentParseError(Exception):
    """文件解析錯誤"""
    pass


class WordDocumentParser:
    """Word 文件解析器"""
    
    def __init__(self, format_handler: FormatHandler, logger: Optional[logging.Logger] = None):
        """
        初始化 Word 文件解析器
        
        Args:
            format_handler: 格式處理器
            logger: 日誌記錄器
        """
        self.format_handler = format_handler
        self.logger = logger or logging.getLogger(__name__)
        self.number_pattern = re.compile(r'^(\d+)\.\s*(.*)')
    
    def parse_document(self, file_path: str) -> Dict[str, Any]:
        """
        解析 Word 文檔內容
        
        Args:
            file_path: Word 文檔路徑
            
        Returns:
            Dict: 包含文檔內容的字典
            
        Raises:
            DocumentParseError: 解析失敗時拋出
        """
        if not self._validate_file(file_path, '.docx'):
            raise DocumentParseError(f"無效的 Word 文檔: {file_path}")
        
        try:
            doc = Document(file_path)
            
            # 解析基本內容
            basic_content = self._extract_basic_content(doc)
            
            # 解析編號段落
            sections = self._parse_numbered_sections(doc)
            
            # 提取元數據
            metadata = self._extract_metadata(doc)
            
            return {
                'file_path': file_path,
                'basic_content': basic_content,
                'sections': sections,
                'metadata': metadata,
                'total_sections': len(sections),
                'success': True,
                'error': None
            }
            
        except Exception as e:
            self.logger.error(f"解析 Word 文檔失敗: {e}")
            raise DocumentParseError(f"解析 Word 文檔失敗: {e}")
    
    def parse_numbered_sections(self, file_path: str) -> Dict[str, Any]:
        """
        解析編號段落（向後兼容的方法）
        
        Args:
            file_path: Word 文檔路徑
            
        Returns:
            Dict: 解析結果
        """
        try:
            result = self.parse_document(file_path)
            return {
                'sections': result['sections'],
                'total_sections': result['total_sections'],
                'success': True,
                'error': None
            }
        except DocumentParseError as e:
            return {
                'sections': [],
                'total_sections': 0,
                'success': False,
                'error': str(e)
            }
    
    def _validate_file(self, file_path: str, extension: str) -> bool:
        """驗證檔案"""
        return (os.path.exists(file_path) and 
                file_path.lower().endswith(extension))
    
    def _extract_basic_content(self, doc: Document) -> Dict[str, Any]:
        """提取基本內容"""
        paragraphs = []
        tables = []
        full_text = []
        
        # 提取段落
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                paragraph_data = {
                    'text': paragraph.text,
                    'style': paragraph.style.name if paragraph.style else None,
                    'formatting': self.format_handler.extract_word_formatting(paragraph)
                }
                paragraphs.append(paragraph_data)
                full_text.append(paragraph.text)
        
        # 提取表格
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)
        
        return {
            'paragraphs': paragraphs,
            'tables': tables,
            'text': '\n'.join(full_text)
        }
    
    def _parse_numbered_sections(self, doc: Document) -> List[Dict[str, Any]]:
        """解析編號段落"""
        sections = []
        current_section = None
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            
            # 提取格式信息
            paragraph_formatting = self.format_handler.extract_word_formatting(paragraph)
            
            match = self.number_pattern.match(text)
            if match:
                # 保存前一個段落
                if current_section is not None:
                    sections.append(current_section)
                
                # 開始新段落
                section_number = int(match.group(1))
                section_title = match.group(2) if match.group(2) else ""
                
                current_section = {
                    'number': section_number,
                    'title': section_title,
                    'content': [text],
                    'text_only': section_title,
                    'formatting': [paragraph_formatting]
                }
            else:
                if current_section is not None:
                    current_section['content'].append(text)
                    current_section['formatting'].append(paragraph_formatting)
                    if current_section['text_only']:
                        current_section['text_only'] += '\n' + text
                    else:
                        current_section['text_only'] = text
                else:
                    # 創建前言段落
                    current_section = {
                        'number': 0,
                        'title': '前言',
                        'content': [text],
                        'text_only': text,
                        'formatting': [paragraph_formatting]
                    }
        
        # 保存最後一個段落
        if current_section is not None:
            sections.append(current_section)
        
        return sections
    
    def _extract_metadata(self, doc: Document) -> Dict[str, Any]:
        """提取元數據"""
        try:
            core_props = doc.core_properties
            return {
                'title': core_props.title,
                'author': core_props.author,
                'subject': core_props.subject,
                'created': core_props.created,
                'modified': core_props.modified,
                'category': core_props.category,
                'comments': core_props.comments
            }
        except Exception as e:
            self.logger.warning(f"提取元數據失敗: {e}")
            return {}


class PowerPointDocumentParser:
    """PowerPoint 文件解析器"""
    
    def __init__(self, format_handler: FormatHandler, logger: Optional[logging.Logger] = None):
        """
        初始化 PowerPoint 文件解析器
        
        Args:
            format_handler: 格式處理器
            logger: 日誌記錄器
        """
        self.format_handler = format_handler
        self.logger = logger or logging.getLogger(__name__)
    
    def parse_document(self, file_path: str) -> Dict[str, Any]:
        """
        解析 PowerPoint 文檔內容
        
        Args:
            file_path: PowerPoint 文檔路徑
            
        Returns:
            Dict: 包含文檔內容的字典
            
        Raises:
            DocumentParseError: 解析失敗時拋出
        """
        if not self._validate_file(file_path, '.pptx'):
            raise DocumentParseError(f"無效的 PowerPoint 文檔: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            slides = self._parse_slides(prs)
            
            return {
                'file_path': file_path,
                'slides': slides,
                'total_slides': len(slides),
                'text': self._extract_all_text(slides),
                'success': True,
                'error': None
            }
            
        except Exception as e:
            self.logger.error(f"解析 PowerPoint 文檔失敗: {e}")
            raise DocumentParseError(f"解析 PowerPoint 文檔失敗: {e}")
    
    def analyze_template_slide(self, slide) -> Dict[str, Any]:
        """
        分析模板投影片的格式特徵
        
        Args:
            slide: 投影片物件
            
        Returns:
            Dict: 分析結果
        """
        analysis = {
            'text_shapes': [],
            'background_shapes': [],
            'has_background_image': False,
            'layout_name': '',
            'master_name': '',
            'summary': ''
        }
        
        try:
            # 分析佈局信息
            if hasattr(slide, 'slide_layout'):
                analysis['layout_name'] = slide.slide_layout.name
                if hasattr(slide.slide_layout, 'slide_master'):
                    analysis['master_name'] = getattr(slide.slide_layout.slide_master, 'name', 'Unknown')
            
            # 分析形狀
            text_count, image_count, other_count = 0, 0, 0
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and hasattr(shape, 'text'):
                    text_count += 1
                    analysis['text_shapes'].append({
                        'type': 'text',
                        'has_text': bool(shape.text.strip()),
                        'width': shape.width,
                        'height': shape.height,
                        'left': shape.left,
                        'top': shape.top
                    })
                elif hasattr(shape, 'shape_type') and shape.shape_type.name == 'PICTURE':
                    image_count += 1
                    # 檢查是否可能是背景圖片
                    is_background = self._is_background_image(shape, slide)
                    if is_background:
                        analysis['has_background_image'] = True
                    
                    analysis['background_shapes'].append({
                        'type': 'picture',
                        'is_background': is_background,
                        'width': shape.width,
                        'height': shape.height
                    })
                else:
                    other_count += 1
            
            # 生成摘要
            summary_parts = []
            if text_count > 0:
                summary_parts.append(f"{text_count}個文本框")
            if image_count > 0:
                summary_parts.append(f"{image_count}個圖片")
            if other_count > 0:
                summary_parts.append(f"{other_count}個其他形狀")
            if analysis['has_background_image']:
                summary_parts.append("有背景圖片")
            
            analysis['summary'] = ', '.join(summary_parts) if summary_parts else '無特殊元素'
            
        except Exception as e:
            self.logger.error(f"分析模板投影片失敗: {e}")
            analysis['summary'] = f"分析失敗: {str(e)}"
        
        return analysis
    
    def _validate_file(self, file_path: str, extension: str) -> bool:
        """驗證檔案"""
        return (os.path.exists(file_path) and 
                file_path.lower().endswith(extension))
    
    def _parse_slides(self, prs: Presentation) -> List[Dict[str, Any]]:
        """解析投影片"""
        slides = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                'slide_number': slide_num,
                'title': '',
                'content': [],
                'text_runs': [],
                'layout_name': slide.slide_layout.name if slide.slide_layout else 'Unknown'
            }
            
            slide_text = self._extract_slide_text(slide, slide_data)
            slide_data['content'] = slide_text
            slide_data['full_text'] = '\n'.join(slide_text)
            
            # 如果沒有識別到標題，使用第一行文本作為標題
            if not slide_data['title'] and slide_text:
                slide_data['title'] = slide_text[0][:50] + ('...' if len(slide_text[0]) > 50 else '')
            
            slides.append(slide_data)
        
        return slides
    
    def _extract_slide_text(self, slide, slide_data: Dict[str, Any]) -> List[str]:
        """提取投影片文本"""
        slide_text = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content = shape.text.strip()
                slide_text.append(text_content)
                slide_data['text_runs'].append({
                    'text': text_content,
                    'shape_type': str(shape.shape_type) if hasattr(shape, 'shape_type') else 'Unknown'
                })
                
                # 嘗試識別標題
                if not slide_data['title'] and (
                    len(text_content) < 100 or
                    'title' in text_content.lower()[:20]
                ):
                    slide_data['title'] = text_content
            
            # 處理表格內容
            if hasattr(shape, 'has_table') and shape.has_table:
                table_data = self._extract_table_data(shape.table)
                if table_data:
                    slide_data['text_runs'].append({
                        'text': f"[表格: {len(table_data)}行]",
                        'shape_type': 'Table',
                        'table_data': table_data
                    })
                    slide_text.extend([cell for row in table_data for cell in row if cell])
        
        return slide_text
    
    def _extract_table_data(self, table) -> List[List[str]]:
        """提取表格數據"""
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_data.append(cell_text)
            if row_data:
                table_data.append(row_data)
        return table_data
    
    def _extract_all_text(self, slides: List[Dict[str, Any]]) -> str:
        """提取所有文本"""
        all_text = []
        for slide in slides:
            all_text.extend(slide['content'])
        return '\n'.join(all_text)
    
    def _is_background_image(self, shape, slide) -> bool:
        """判斷是否為背景圖片"""
        try:
            slide_width = slide.slide_layout.width if hasattr(slide.slide_layout, 'width') else 0
            slide_height = slide.slide_layout.height if hasattr(slide.slide_layout, 'height') else 0
            
            return (slide_width > 0 and slide_height > 0 and
                    shape.width >= slide_width * 0.8 and 
                    shape.height >= slide_height * 0.8)
        except:
            return False


class DocumentParserFactory:
    """文件解析器工廠"""
    
    @staticmethod
    def create_parser(file_type: str, format_handler: FormatHandler, 
                     logger: Optional[logging.Logger] = None):
        """
        創建文件解析器
        
        Args:
            file_type: 文件類型 ('word' 或 'powerpoint')
            format_handler: 格式處理器
            logger: 日誌記錄器
            
        Returns:
            文件解析器實例
            
        Raises:
            ValueError: 不支持的文件類型
        """
        if file_type.lower() in ['word', 'docx']:
            return WordDocumentParser(format_handler, logger)
        elif file_type.lower() in ['powerpoint', 'pptx']:
            return PowerPointDocumentParser(format_handler, logger)
        else:
            raise ValueError(f"不支持的文件類型: {file_type}")
    
    @staticmethod
    def create_parser_from_file(file_path: str, format_handler: FormatHandler,
                               logger: Optional[logging.Logger] = None):
        """
        根據檔案路徑創建解析器
        
        Args:
            file_path: 檔案路徑
            format_handler: 格式處理器
            logger: 日誌記錄器
            
        Returns:
            文件解析器實例
        """
        if file_path.lower().endswith('.docx'):
            return WordDocumentParser(format_handler, logger)
        elif file_path.lower().endswith('.pptx'):
            return PowerPointDocumentParser(format_handler, logger)
        else:
            raise ValueError(f"不支持的檔案格式: {file_path}")