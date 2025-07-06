"""
文件轉換器主模組 - 整合所有功能的主要轉換器
使用策略模式和工廠模式實現靈活的文件轉換
"""

from typing import Dict, List, Any, Optional, Callable
from abc import ABC, abstractmethod
from pptx import Presentation
import os
from pathlib import Path

# 導入我們的模組
from format_handler import FormatHandler
from document_parser import DocumentParserFactory, WordDocumentParser, PowerPointDocumentParser
from slide_manager import SlideManager, SlideAnalyzer
from logger_config import (
    LoggerConfig, ErrorHandler, PerformanceMonitor, 
    ConversionError, DocumentError, create_result_dict, get_logger
)


class ConversionStrategy(ABC):
    """轉換策略抽象基類"""
    
    @abstractmethod
    def convert(self, source_file: str, template_file: str, output_file: Optional[str] = None,
               progress_callback: Optional[Callable] = None) -> Dict[str, Any]:
        """執行轉換"""
        pass


class WordToPowerPointStrategy(ConversionStrategy):
    """Word 轉 PowerPoint 策略"""
    
    def __init__(self, format_handler: FormatHandler, slide_manager: SlideManager,
                 word_parser: WordDocumentParser, ppt_parser: PowerPointDocumentParser,
                 error_handler: ErrorHandler, performance_monitor: PerformanceMonitor):
        """
        初始化 Word 轉 PowerPoint 策略
        
        Args:
            format_handler: 格式處理器
            slide_manager: 投影片管理器
            word_parser: Word 解析器
            ppt_parser: PowerPoint 解析器
            error_handler: 錯誤處理器
            performance_monitor: 性能監控器
        """
        self.format_handler = format_handler
        self.slide_manager = slide_manager
        self.word_parser = word_parser
        self.ppt_parser = ppt_parser
        self.error_handler = error_handler
        self.performance_monitor = performance_monitor
        self.logger = get_logger(self.__class__.__name__)
    
    def convert(self, source_file: str, template_file: str, output_file: Optional[str] = None,
               progress_callback: Optional[Callable] = None) -> Dict[str, Any]:
        """
        執行 Word 轉 PowerPoint 轉換
        
        Args:
            source_file: Word 檔案路徑
            template_file: PowerPoint 模板路徑
            output_file: 輸出檔案路徑
            progress_callback: 進度回調函數 (current, total, message)
            
        Returns:
            Dict: 轉換結果
        """
        operation_name = "Word轉PowerPoint"
        self.performance_monitor.start_timing(operation_name)
        
        try:
            self.error_handler.log_operation_start(operation_name, {
                'source': source_file,
                'template': template_file,
                'output': output_file
            })
            
            # 1. 解析 Word 文檔
            if progress_callback:
                progress_callback(1, 5, "解析 Word 文檔...")
            
            try:
                word_data = self.word_parser.parse_document(source_file)
                sections = word_data['sections']
                self.logger.info(f"成功解析 Word 文檔，找到 {len(sections)} 個段落")
            except Exception as e:
                raise DocumentError(f"解析 Word 文檔失敗: {str(e)}", "WORD_PARSE_ERROR")
            
            # 2. 載入 PowerPoint 模板
            if progress_callback:
                progress_callback(2, 5, "載入 PowerPoint 模板...")
            
            try:
                if not os.path.exists(template_file):
                    raise DocumentError(f"模板檔案不存在: {template_file}", "TEMPLATE_NOT_FOUND")
                
                prs = Presentation(template_file)
                if len(prs.slides) == 0:
                    raise DocumentError("模板中沒有投影片", "EMPTY_TEMPLATE")
                
                template_slide = prs.slides[0]
                self.logger.info(f"成功載入模板，包含 {len(prs.slides)} 張投影片")
            except DocumentError:
                raise
            except Exception as e:
                raise DocumentError(f"載入模板失敗: {str(e)}", "TEMPLATE_LOAD_ERROR")
            
            # 3. 分析模板
            if progress_callback:
                progress_callback(3, 5, "分析模板結構...")
            
            template_analysis = self.ppt_parser.analyze_template_slide(template_slide)
            self.logger.debug(f"模板分析結果: {template_analysis['summary']}")
            
            # 4. 執行轉換
            if progress_callback:
                progress_callback(4, 5, "轉換內容...")
            
            def slide_progress(current, total, message):
                overall_progress = 4 + (current / total) * 0.8  # 4-4.8 的範圍
                if progress_callback:
                    progress_callback(overall_progress, 5, f"轉換投影片: {message}")
            
            conversion_result = self.slide_manager.replace_slides_with_sections(
                prs, sections, template_slide, slide_progress)
            
            if not conversion_result['success']:
                raise ConversionError(f"投影片轉換失敗: {conversion_result.get('error', '未知錯誤')}")
            
            # 5. 保存檔案
            if progress_callback:
                progress_callback(5, 5, "保存檔案...")
            
            if output_file is None:
                output_file = self._generate_output_filename(template_file)
            
            # 確保輸出目錄存在
            output_path = Path(output_file)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            prs.save(output_file)
            
            # 記錄成功
            duration = self.performance_monitor.end_timing(operation_name)
            
            result = create_result_dict(
                success=True,
                total_sections=len(sections),
                slides_created=conversion_result['slides_created'],
                output_file=output_file,
                skipped_sections=conversion_result.get('skipped_sections', []),
                format_issues=conversion_result.get('format_issues', []),
                processing_time=duration,
                template_analysis=template_analysis
            )
            
            self.error_handler.log_operation_success(operation_name, result)
            
            # 記錄警告（如果有）
            if result['skipped_sections']:
                self.error_handler.log_operation_warning(
                    operation_name, 
                    f"跳過了 {len(result['skipped_sections'])} 個段落",
                    {'skipped_count': len(result['skipped_sections'])}
                )
            
            return result
            
        except Exception as e:
            duration = self.performance_monitor.end_timing(operation_name)
            error_info = self.error_handler.handle_error(e, operation_name)
            
            return create_result_dict(
                success=False,
                error=str(e),
                error_info=error_info,
                processing_time=duration
            )
    
    def _generate_output_filename(self, template_file: str) -> str:
        """生成輸出檔案名"""
        template_path = Path(template_file)
        return str(template_path.parent / f"{template_path.stem}_轉換版{template_path.suffix}")


class DocumentConverter:
    """主要文件轉換器 - 使用策略模式"""
    
    def __init__(self, strategy: Optional[ConversionStrategy] = None, 
                 logger_level: str = "INFO", log_to_file: bool = True):
        """
        初始化文件轉換器
        
        Args:
            strategy: 轉換策略
            logger_level: 日誌級別
            log_to_file: 是否記錄到檔案
        """
        # 設置日誌
        from logger_config import LogLevel
        level_map = {
            "DEBUG": LogLevel.DEBUG,
            "INFO": LogLevel.INFO,
            "WARNING": LogLevel.WARNING,
            "ERROR": LogLevel.ERROR
        }
        log_level = level_map.get(logger_level.upper(), LogLevel.INFO)
        
        self.logger = LoggerConfig.setup_logger(
            "DocumentConverter", 
            level=log_level, 
            file_output=log_to_file
        )
        
        # 初始化組件
        self.format_handler = FormatHandler(self.logger)
        self.error_handler = ErrorHandler(self.logger)
        self.performance_monitor = PerformanceMonitor(self.logger)
        
        # 解析器
        self.word_parser = DocumentParserFactory.create_parser('word', self.format_handler, self.logger)
        self.ppt_parser = DocumentParserFactory.create_parser('powerpoint', self.format_handler, self.logger)
        
        # 投影片管理器
        self.slide_manager = SlideManager(self.format_handler, self.logger)
        self.slide_analyzer = SlideAnalyzer(self.logger)
        
        # 設置策略
        self._strategy = strategy or self._create_default_strategy()
    
    def set_strategy(self, strategy: ConversionStrategy):
        """設置轉換策略"""
        self._strategy = strategy
        self.logger.info(f"轉換策略已更改為: {strategy.__class__.__name__}")
    
    def convert_document(self, source_file: str, template_file: str, 
                        output_file: Optional[str] = None,
                        progress_callback: Optional[Callable] = None) -> Dict[str, Any]:
        """
        轉換文檔
        
        Args:
            source_file: 源文件路徑
            template_file: 模板文件路徑
            output_file: 輸出文件路徑
            progress_callback: 進度回調函數
            
        Returns:
            Dict: 轉換結果
        """
        try:
            # 驗證檔案
            self._validate_files(source_file, template_file)
            
            # 執行轉換
            return self._strategy.convert(source_file, template_file, output_file, progress_callback)
            
        except Exception as e:
            error_info = self.error_handler.handle_error(e, "文檔轉換")
            return create_result_dict(success=False, error=str(e), error_info=error_info)
    
    def analyze_document(self, file_path: str) -> Dict[str, Any]:
        """
        分析文檔結構
        
        Args:
            file_path: 文檔路徑
            
        Returns:
            Dict: 分析結果
        """
        try:
            if file_path.lower().endswith('.docx'):
                return self.word_parser.parse_document(file_path)
            elif file_path.lower().endswith('.pptx'):
                ppt_data = self.ppt_parser.parse_document(file_path)
                # 添加結構分析
                prs = Presentation(file_path)
                structure_analysis = self.slide_analyzer.analyze_presentation_structure(prs)
                ppt_data['structure_analysis'] = structure_analysis
                return ppt_data
            else:
                raise ValueError(f"不支持的檔案格式: {file_path}")
                
        except Exception as e:
            error_info = self.error_handler.handle_error(e, "文檔分析")
            return create_result_dict(success=False, error=str(e), error_info=error_info)
    
    def get_conversion_preview(self, source_file: str) -> Dict[str, Any]:
        """
        獲取轉換預覽信息
        
        Args:
            source_file: 源文件路徑
            
        Returns:
            Dict: 預覽信息
        """
        try:
            if not source_file.lower().endswith('.docx'):
                raise ValueError("只支持 Word 文檔預覽")
            
            word_data = self.word_parser.parse_document(source_file)
            sections = word_data['sections']
            
            preview = {
                'success': True,
                'total_sections': len(sections),
                'estimated_slides': len(sections),
                'sections_preview': []
            }
            
            for section in sections[:5]:  # 只顯示前5個段落的預覽
                preview['sections_preview'].append({
                    'number': section['number'],
                    'title': section['title'][:50] + ('...' if len(section['title']) > 50 else ''),
                    'content_length': len(section['text_only']),
                    'has_formatting': bool(section.get('formatting'))
                })
            
            if len(sections) > 5:
                preview['sections_preview'].append({
                    'number': '...',
                    'title': f'還有 {len(sections) - 5} 個段落',
                    'content_length': 0,
                    'has_formatting': False
                })
            
            return preview
            
        except Exception as e:
            error_info = self.error_handler.handle_error(e, "轉換預覽")
            return create_result_dict(success=False, error=str(e), error_info=error_info)
    
    def _create_default_strategy(self) -> ConversionStrategy:
        """創建預設轉換策略"""
        return WordToPowerPointStrategy(
            self.format_handler,
            self.slide_manager,
            self.word_parser,
            self.ppt_parser,
            self.error_handler,
            self.performance_monitor
        )
    
    def _validate_files(self, source_file: str, template_file: str):
        """驗證檔案"""
        if not os.path.exists(source_file):
            raise DocumentError(f"源文件不存在: {source_file}", "SOURCE_NOT_FOUND")
        
        if not os.path.exists(template_file):
            raise DocumentError(f"模板文件不存在: {template_file}", "TEMPLATE_NOT_FOUND")
        
        if not source_file.lower().endswith('.docx'):
            raise DocumentError(f"不支援的源文件格式: {source_file}", "UNSUPPORTED_SOURCE_FORMAT")
        
        if not template_file.lower().endswith('.pptx'):
            raise DocumentError(f"不支援的模板格式: {template_file}", "UNSUPPORTED_TEMPLATE_FORMAT")


class ConverterFactory:
    """轉換器工廠"""
    
    @staticmethod
    def create_converter(converter_type: str = "word_to_ppt", **kwargs) -> DocumentConverter:
        """
        創建轉換器
        
        Args:
            converter_type: 轉換器類型
            **kwargs: 其他參數
            
        Returns:
            DocumentConverter: 轉換器實例
        """
        if converter_type.lower() == "word_to_ppt":
            return DocumentConverter(**kwargs)
        else:
            raise ValueError(f"不支持的轉換器類型: {converter_type}")
    
    @staticmethod
    def create_batch_converter(**kwargs) -> 'BatchConverter':
        """創建批次轉換器"""
        return BatchConverter(**kwargs)


class BatchConverter:
    """批次轉換器"""
    
    def __init__(self, **kwargs):
        """初始化批次轉換器"""
        self.converter = DocumentConverter(**kwargs)
        self.logger = self.converter.logger
    
    def convert_multiple(self, file_pairs: List[Dict[str, str]], 
                        progress_callback: Optional[Callable] = None) -> List[Dict[str, Any]]:
        """
        批次轉換多個文檔
        
        Args:
            file_pairs: 文件對列表，每個包含 source, template, output
            progress_callback: 進度回調函數
            
        Returns:
            List[Dict]: 轉換結果列表
        """
        results = []
        total = len(file_pairs)
        
        for i, file_pair in enumerate(file_pairs):
            try:
                if progress_callback:
                    progress_callback(i + 1, total, f"轉換檔案 {i + 1}/{total}")
                
                result = self.converter.convert_document(
                    file_pair['source'],
                    file_pair['template'],
                    file_pair.get('output')
                )
                
                result['source_file'] = file_pair['source']
                results.append(result)
                
            except Exception as e:
                error_result = create_result_dict(
                    success=False,
                    error=str(e),
                    source_file=file_pair['source']
                )
                results.append(error_result)
                self.logger.error(f"批次轉換失敗: {file_pair['source']} - {e}")
        
        successful = sum(1 for r in results if r['success'])
        self.logger.info(f"批次轉換完成: {successful}/{total} 成功")
        
        return results


# 便利函數
def convert_word_to_ppt(word_file: str, ppt_template: str, output_file: Optional[str] = None,
                       progress_callback: Optional[Callable] = None) -> Dict[str, Any]:
    """
    便利函數：Word 轉 PowerPoint
    
    Args:
        word_file: Word 檔案路徑
        ppt_template: PowerPoint 模板路徑
        output_file: 輸出檔案路徑
        progress_callback: 進度回調函數
        
    Returns:
        Dict: 轉換結果
    """
    converter = ConverterFactory.create_converter()
    return converter.convert_document(word_file, ppt_template, output_file, progress_callback)


def analyze_document_structure(file_path: str) -> Dict[str, Any]:
    """
    便利函數：分析文檔結構
    
    Args:
        file_path: 檔案路徑
        
    Returns:
        Dict: 分析結果
    """
    converter = ConverterFactory.create_converter()
    return converter.analyze_document(file_path)