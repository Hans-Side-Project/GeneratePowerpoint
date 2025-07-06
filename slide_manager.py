"""
投影片管理模組 - 負責 PowerPoint 投影片的操作和管理
提供投影片複製、內容替換、格式處理等功能
"""

from typing import Dict, List, Any, Optional, Callable
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
import os
import io
import logging
from format_handler import FormatHandler


class SlideOperationError(Exception):
    """投影片操作錯誤"""
    pass


class SlideManager:
    """投影片管理器"""
    
    def __init__(self, format_handler: FormatHandler, logger: Optional[logging.Logger] = None):
        """
        初始化投影片管理器
        
        Args:
            format_handler: 格式處理器
            logger: 日誌記錄器
        """
        self.format_handler = format_handler
        self.logger = logger or logging.getLogger(__name__)
    
    def duplicate_slide(self, ppt_file_path: str, source_slide_number: int, 
                       copy_count: int = 1, output_file: Optional[str] = None) -> Dict[str, Any]:
        """
        複製指定的投影片
        
        Args:
            ppt_file_path: PowerPoint 檔案路徑
            source_slide_number: 源投影片編號（從1開始）
            copy_count: 複製數量
            output_file: 輸出檔案路徑
            
        Returns:
            Dict: 操作結果
            
        Raises:
            SlideOperationError: 操作失敗時拋出
        """
        try:
            if not os.path.exists(ppt_file_path):
                raise SlideOperationError(f"檔案不存在: {ppt_file_path}")
            
            prs = Presentation(ppt_file_path)
            
            if source_slide_number < 1 or source_slide_number > len(prs.slides):
                raise SlideOperationError(f"投影片編號無效: {source_slide_number}")
            
            # 獲取源投影片
            source_slide = prs.slides[source_slide_number - 1]
            source_layout = source_slide.slide_layout
            
            # 複製指定次數
            for i in range(copy_count):
                new_slide = prs.slides.add_slide(source_layout)
                self._copy_slide_completely(source_slide, new_slide)
            
            # 保存檔案
            save_path = output_file or ppt_file_path
            prs.save(save_path)
            
            result = {
                'success': True,
                'copied_slides': copy_count,
                'total_slides_after': len(prs.slides),
                'output_file': save_path,
                'error': None
            }
            
            self.logger.info(f"成功複製 {copy_count} 張投影片")
            return result
            
        except Exception as e:
            error_msg = f"複製投影片失敗: {str(e)}"
            self.logger.error(error_msg)
            raise SlideOperationError(error_msg)
    
    def replace_slides_with_sections(self, prs: Presentation, sections: List[Dict[str, Any]], 
                                   template_slide, progress_callback: Optional[Callable] = None) -> Dict[str, Any]:
        """
        用章節內容替換投影片
        
        Args:
            prs: PowerPoint 演示文稿物件
            sections: 章節資料列表
            template_slide: 模板投影片
            progress_callback: 進度回調函數
            
        Returns:
            Dict: 操作結果
        """
        result = {
            'success': False,
            'slides_created': 0,
            'skipped_sections': [],
            'format_issues': [],
            'error': None
        }
        
        try:
            template_layout = template_slide.slide_layout
            
            # 清除現有投影片（保留第一張）
            self._clear_existing_slides(prs)
            
            slides_created = 0
            
            for i, section in enumerate(sections):
                try:
                    if progress_callback:
                        progress_callback(i + 1, len(sections), f"處理段落 {section['number']}")
                    
                    if i == 0:
                        target_slide = prs.slides[0]  # 使用第一張投影片
                    else:
                        target_slide = prs.slides.add_slide(template_layout)
                        self._copy_slide_completely(template_slide, target_slide)
                    
                    # 替換內容
                    self._replace_slide_content(target_slide, section, template_slide)
                    slides_created += 1
                    
                    self.logger.debug(f"成功處理段落 {section['number']}")
                    
                except Exception as e:
                    error_info = {
                        'number': section['number'],
                        'title': section.get('title', ''),
                        'error': str(e)
                    }
                    result['skipped_sections'].append(error_info)
                    self.logger.warning(f"跳過段落 {section['number']}: {e}")
            
            result.update({
                'success': True,
                'slides_created': slides_created
            })
            
            self.logger.info(f"成功創建 {slides_created} 張投影片")
            
        except Exception as e:
            error_msg = f"替換投影片內容失敗: {str(e)}"
            result['error'] = error_msg
            self.logger.error(error_msg)
        
        return result
    
    def copy_slide_background(self, source_slide, target_slide) -> bool:
        """
        複製投影片背景
        
        Args:
            source_slide: 源投影片
            target_slide: 目標投影片
            
        Returns:
            bool: 是否成功複製
        """
        try:
            # 複製背景填充
            if hasattr(source_slide, 'background') and hasattr(target_slide, 'background'):
                source_bg = source_slide.background
                target_bg = target_slide.background
                
                if hasattr(source_bg, 'fill') and hasattr(target_bg, 'fill'):
                    # 這裡可以擴展更複雜的背景複製邏輯
                    pass
            
            # 確保使用相同的佈局
            if (hasattr(source_slide, 'slide_layout') and 
                hasattr(target_slide, 'slide_layout')):
                source_layout = source_slide.slide_layout
                target_layout = target_slide.slide_layout
                
                if source_layout.name != target_layout.name:
                    self.logger.warning(f"佈局差異: 源={source_layout.name}, 目標={target_layout.name}")
            
            return True
            
        except Exception as e:
            self.logger.error(f"複製背景失敗: {e}")
            return False
    
    def _copy_slide_completely(self, source_slide, target_slide):
        """完整複製投影片內容"""
        try:
            # 複製背景
            self.copy_slide_background(source_slide, target_slide)
            
            # 複製所有形狀
            for shape in source_slide.shapes:
                try:
                    if shape.is_placeholder:
                        self._copy_placeholder_content(shape, target_slide)
                    else:
                        self._copy_non_placeholder_shape(shape, target_slide)
                except Exception as e:
                    self.logger.warning(f"複製形狀失敗: {e}")
                    continue
                    
        except Exception as e:
            self.logger.error(f"完整複製投影片失敗: {e}")
            raise SlideOperationError(f"複製投影片失敗: {e}")
    
    def _copy_placeholder_content(self, source_placeholder, target_slide):
        """複製占位符內容"""
        try:
            for target_shape in target_slide.shapes:
                if (target_shape.is_placeholder and 
                    hasattr(source_placeholder, 'placeholder_format') and
                    hasattr(target_shape, 'placeholder_format') and
                    target_shape.placeholder_format.idx == source_placeholder.placeholder_format.idx):
                    
                    # 複製文本內容
                    if hasattr(source_placeholder, 'text') and hasattr(target_shape, 'text'):
                        target_shape.text = source_placeholder.text
                    
                    # 複製文本框架內容
                    if (hasattr(source_placeholder, 'text_frame') and 
                        hasattr(target_shape, 'text_frame')):
                        self.format_handler.copy_text_frame(
                            source_placeholder.text_frame, target_shape.text_frame)
                    
                    break
                    
        except Exception as e:
            self.logger.warning(f"複製占位符內容失敗: {e}")
    
    def _copy_non_placeholder_shape(self, source_shape, target_slide):
        """複製非占位符形狀"""
        try:
            left, top, width, height = source_shape.left, source_shape.top, source_shape.width, source_shape.height
            
            if source_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                new_textbox = target_slide.shapes.add_textbox(left, top, width, height)
                if hasattr(source_shape, 'text_frame'):
                    self.format_handler.copy_text_frame(source_shape.text_frame, new_textbox.text_frame)
                self.format_handler.copy_shape_formatting(source_shape, new_textbox)
                
            elif source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_blob = source_shape.image.blob
                    image_stream = io.BytesIO(image_blob)
                    new_picture = target_slide.shapes.add_picture(image_stream, left, top, width, height)
                    self.format_handler.copy_shape_formatting(source_shape, new_picture)
                except Exception as e:
                    self.logger.warning(f"複製圖片失敗: {e}")
                    
            elif source_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    if hasattr(source_shape, 'auto_shape_type'):
                        new_shape = target_slide.shapes.add_shape(
                            source_shape.auto_shape_type, left, top, width, height)
                        if hasattr(source_shape, 'text_frame'):
                            self.format_handler.copy_text_frame(source_shape.text_frame, new_shape.text_frame)
                        self.format_handler.copy_shape_formatting(source_shape, new_shape)
                except Exception as e:
                    self.logger.warning(f"複製自動形狀失敗: {e}")
                    
            elif source_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                try:
                    if hasattr(source_shape, 'table'):
                        table = source_shape.table
                        rows, cols = len(table.rows), len(table.columns)
                        new_table = target_slide.shapes.add_table(rows, cols, left, top, width, height)
                        self._copy_table_content(table, new_table)
                except Exception as e:
                    self.logger.warning(f"複製表格失敗: {e}")
            
            else:
                self.logger.debug(f"跳過未處理的形狀類型: {source_shape.shape_type}")
                
        except Exception as e:
            self.logger.warning(f"複製非占位符形狀失敗: {e}")
    
    def _copy_table_content(self, source_table, target_table_shape):
        """複製表格內容"""
        try:
            target_table = target_table_shape.table if hasattr(target_table_shape, 'table') else target_table_shape
            
            for row_idx, source_row in enumerate(source_table.rows):
                if row_idx < len(target_table.rows):
                    target_row = target_table.rows[row_idx]
                    for col_idx, source_cell in enumerate(source_row.cells):
                        if col_idx < len(target_row.cells):
                            target_cell = target_row.cells[col_idx]
                            target_cell.text = source_cell.text
                            
                            # 複製單元格文本格式
                            if (hasattr(source_cell, 'text_frame') and 
                                hasattr(target_cell, 'text_frame')):
                                self.format_handler.copy_text_frame(
                                    source_cell.text_frame, target_cell.text_frame)
                                    
        except Exception as e:
            self.logger.warning(f"複製表格內容失敗: {e}")
    
    def _replace_slide_content(self, slide, section: Dict[str, Any], template_slide=None):
        """替換投影片內容"""
        try:
            # 查找文本框
            text_shapes = [shape for shape in slide.shapes 
                          if hasattr(shape, 'text_frame') and hasattr(shape, 'text')]
            
            # 如果沒有文本框，從模板創建
            if not text_shapes and template_slide:
                text_shapes = self._create_text_box_from_template(slide, template_slide)
            
            if not text_shapes:
                # 創建默認文本框
                new_textbox = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1), Inches(9), Inches(6.5))
                text_shapes.append(new_textbox)
            
            # 主文本框
            main_text_shape = text_shapes[0]
            
            # 如果有格式化數據，使用詳細格式替換
            if 'formatting' in section and section['formatting']:
                self._replace_content_with_formatting(main_text_shape, section)
            else:
                self._replace_content_basic(main_text_shape, section)
            
            # 清空其他文本框
            for i in range(1, len(text_shapes)):
                text_shapes[i].text = ""
                
        except Exception as e:
            self.logger.error(f"替換投影片內容失敗: {e}")
            raise SlideOperationError(f"替換內容失敗: {e}")
    
    def _create_text_box_from_template(self, slide, template_slide) -> List:
        """從模板創建文本框"""
        text_shapes = []
        template_text_shapes = [shape for shape in template_slide.shapes 
                               if hasattr(shape, 'text_frame') and hasattr(shape, 'text')]
        
        if template_text_shapes:
            template_shape = template_text_shapes[0]
            new_textbox = slide.shapes.add_textbox(
                template_shape.left, template_shape.top, 
                template_shape.width, template_shape.height)
            
            # 複製格式
            self.format_handler.copy_shape_formatting(template_shape, new_textbox)
            
            # 複製文本框架格式
            if hasattr(template_shape, 'text_frame') and hasattr(new_textbox, 'text_frame'):
                self.format_handler._copy_text_frame_formatting(
                    template_shape.text_frame, new_textbox.text_frame)
            
            text_shapes.append(new_textbox)
            self.logger.debug("成功從模板創建文本框")
        
        return text_shapes
    
    def _replace_content_with_formatting(self, main_text_shape, section: Dict[str, Any]):
        """使用格式化數據替換內容"""
        try:
            main_text_shape.text_frame.clear()
            
            for content_idx, (content_line, formatting_data) in enumerate(
                zip(section['content'], section.get('formatting', []))):
                
                # 創建段落
                if content_idx == 0:
                    paragraph = main_text_shape.text_frame.paragraphs[0]
                else:
                    paragraph = main_text_shape.text_frame.add_paragraph()
                
                # 為每個格式化的文本片段創建運行
                for format_info in formatting_data:
                    if format_info['text'].strip():
                        run = paragraph.add_run()
                        run.text = format_info['text']
                        self.format_handler.apply_word_formatting_to_ppt(run, format_info)
                        
        except Exception as e:
            self.logger.warning(f"格式化替換失敗，使用基本替換: {e}")
            self._replace_content_basic(main_text_shape, section)
    
    def _replace_content_basic(self, main_text_shape, section: Dict[str, Any]):
        """基本內容替換"""
        try:
            # 準備內容
            if section['number'] == 0:
                title_text = section['title']
            else:
                title_text = f"{section['number']}. {section['title']}"
            
            content_lines = section['content']
            if len(content_lines) > 1:
                content_text = '\n'.join(content_lines[1:])
            else:
                content_text = section['text_only']
            
            if not content_text.strip():
                content_text = title_text
            
            # 設置文本
            if content_text != title_text and content_text.strip():
                full_text = f"{title_text}\n{content_text}"
            else:
                full_text = title_text
            
            # 過濾空行
            lines = [line for line in full_text.split('\n') if line.strip()]
            main_text_shape.text = '\n'.join(lines)
            
        except Exception as e:
            self.logger.error(f"基本內容替換失敗: {e}")
            main_text_shape.text = f"錯誤: 無法顯示內容"
    
    def _clear_existing_slides(self, prs: Presentation):
        """清除現有投影片（保留第一張）"""
        try:
            for i in range(len(prs.slides) - 1, 0, -1):
                slide_to_remove = prs.slides[i]
                rId = prs.slides._slides[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._slides[i]
            
            self.logger.debug(f"清除了 {len(prs.slides) - 1} 張現有投影片")
            
        except Exception as e:
            self.logger.warning(f"清除現有投影片失敗: {e}")


class SlideAnalyzer:
    """投影片分析器"""
    
    def __init__(self, logger: Optional[logging.Logger] = None):
        """
        初始化投影片分析器
        
        Args:
            logger: 日誌記錄器
        """
        self.logger = logger or logging.getLogger(__name__)
    
    def analyze_presentation_structure(self, prs: Presentation) -> Dict[str, Any]:
        """
        分析演示文稿結構
        
        Args:
            prs: PowerPoint 演示文稿物件
            
        Returns:
            Dict: 分析結果
        """
        analysis = {
            'total_slides': len(prs.slides),
            'layouts_used': set(),
            'text_shapes_count': 0,
            'image_shapes_count': 0,
            'table_shapes_count': 0,
            'has_master_slide': False,
            'slides_analysis': []
        }
        
        try:
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_analysis = self._analyze_single_slide(slide, slide_num)
                analysis['slides_analysis'].append(slide_analysis)
                
                # 統計信息
                analysis['layouts_used'].add(slide_analysis['layout_name'])
                analysis['text_shapes_count'] += slide_analysis['text_shapes']
                analysis['image_shapes_count'] += slide_analysis['image_shapes']
                analysis['table_shapes_count'] += slide_analysis['table_shapes']
            
            analysis['layouts_used'] = list(analysis['layouts_used'])
            
        except Exception as e:
            self.logger.error(f"分析演示文稿結構失敗: {e}")
        
        return analysis
    
    def _analyze_single_slide(self, slide, slide_num: int) -> Dict[str, Any]:
        """分析單張投影片"""
        analysis = {
            'slide_number': slide_num,
            'layout_name': slide.slide_layout.name if slide.slide_layout else 'Unknown',
            'text_shapes': 0,
            'image_shapes': 0,
            'table_shapes': 0,
            'other_shapes': 0,
            'has_title': False,
            'estimated_word_count': 0
        }
        
        word_count = 0
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and hasattr(shape, 'text'):
                analysis['text_shapes'] += 1
                text_content = shape.text.strip()
                if text_content:
                    word_count += len(text_content.split())
                    if not analysis['has_title'] and len(text_content) < 100:
                        analysis['has_title'] = True
                        
            elif hasattr(shape, 'shape_type'):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    analysis['image_shapes'] += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    analysis['table_shapes'] += 1
                else:
                    analysis['other_shapes'] += 1
        
        analysis['estimated_word_count'] = word_count
        return analysis